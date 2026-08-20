"""
VIKA Wealth — Portfolio Stats (single-file build)

This file combines what used to be three separate modules:
  - portfolio_pipeline.py  -> data loading, ledger/TWR math, metrics
  - portfolio_pdf.py       -> PDF report generation (reportlab)
  - app.py                 -> the Streamlit UI itself

They're kept as clearly-labeled sections below (search for "SECTION:") rather
than actually interleaved, so the internal structure is easy to follow.
Splitting into files was not the cause of the data/download bugs fixed
earlier; this merge is purely a deployment-simplicity preference.
"""
from __future__ import annotations

import io
import os as _os
import time
from dataclasses import dataclass
from datetime import date, datetime
from pathlib import Path
from typing import Iterable, Optional

import matplotlib
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
import pandas as pd
import streamlit as st
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import (
    HRFlowable,
    Image,
    KeepTogether,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont


# ============================== SECTION: pipeline (data + math) ==============================

REQUIRED_COLUMNS = {"portfolio", "date", "ticker", "action", "quantity", "price"}
REQUIRED_HOLDINGS_COLUMNS = {"portfolio", "startdate", "ticker", "quantity", "avgcostprice"}
REQUIRED_CASHFLOW_COLUMNS = {"portfolio", "date", "amount"}

# Shared correlation heatmap colormap: a high-contrast diverging scheme —
# deep navy at -1 (strong negative / diversifying), cream at 0 (uncorrelated),
# deep gold/amber at +1 (strong positive / concentrated risk). Diverging
# through a light neutral midpoint (rather than through a mid-brightness
# yellow) makes the strength of correlation visually obvious at a glance:
# saturated color = strong relationship, near-white = none.
# Import matplotlib.colors lazily to avoid hard dependency at module load time.
def get_corr_cmap():
    import matplotlib.colors as mcolors
    return mcolors.LinearSegmentedColormap.from_list(
        "corr_cmap",
        [
            (0.0, "#290084"),   # -1: brand indigo
            (0.25, "#8B6BC0"),  # -0.5: lighter indigo/lavender
            (0.5, "#FBF7EF"),   # 0: cream (near-white, low saturation)
            (0.75, "#F6BF02"),  # +0.5: brand gold
            (1.0, "#8C5A0B"),   # +1: deep gold/brown
        ],
    )

TICKER_MAP: dict[str, str] = {
    # corrections
    "QPOWER": "QPOWER.NS",
    "APARINDS": "APARINDS.NS",
    # explicit mappings
    "SAGILITY": "SAGILITY.NS",
    "TDPOWERSYS": "TDPOWERSYS.NS",
    # REIT fix
    "MINDSPACE-RR": "MINDSPACE-RR.NS",
}

BSE_OVERRIDE: dict[str, str] = {
    "MINDSPACE-RR": "MINDSPACE",
}

NS_TO_BSE_TICKER: dict[str, str] = {
    # Some instruments (especially ETFs) use numeric symbols on BSE.
    # MODEFENCE ETF: BSE symbol 590152 (per issuer factsheet).
    "MODEFENCE.NS": "590152.BO",
}

# Keep non-company instruments out of the misleading Large/Mid/Small buckets.
# ETFs do not have a company-equity market cap. REITs do, but Yahoo can
# occasionally omit it, so REIT is used as the fallback asset class.
ETF_TICKERS: set[str] = {
    "GOLDBEES.NS",
    "SILVERBEES.NS",
    "MODEFENCE.NS",
}

REIT_TICKERS: set[str] = {
    "MINDSPACE-RR.NS",
    "EMBASSY.NS",
    "BROOKFIELD.NS",
    "NEXUS.NS",
}

ASSET_CLASS_MAP: dict[str, str] = {
    **{t: "ETF" for t in ETF_TICKERS},
    **{t: "REIT" for t in REIT_TICKERS},
}

BENCHMARK_NIFTY500_TICKER = "0P0001IAU3.BO"  # Nifty 500 TRI proxy (mutual fund NAV)
BENCHMARK_NIFTY50_TICKER = "0P00005WL6.BO"  # Nifty 50 TRI proxy (mutual fund NAV)
BENCHMARK_TICKER = BENCHMARK_NIFTY500_TICKER  # kept for backward compatibility
BENCHMARKS: dict[str, str] = {"Nifty 50": BENCHMARK_NIFTY50_TICKER, "Nifty 500": BENCHMARK_NIFTY500_TICKER}

SECTOR_MAP: dict[str, str] = {
    # Financials
    "ABCAPITAL.NS": "BFSI",
    "TATACAP.NS": "BFSI",
    "FEDERALBNK.NS": "BFSI",
    "SHRIRAMFIN.NS": "BFSI",
    # Industrials / Capital Goods
    "LT.NS": "Capital Goods",
    "BHEL.NS": "Capital Goods",
    "KPIL.NS": "Capital Goods",
    "ELECON.NS": "Capital Goods",
    "TRITURBINE.NS": "Capital Goods",
    "INOXWIND.NS": "Renewable",
    "GENUSPOWER.NS": "Capital Goods",
    "TDPOWERSYS.NS": "Capital Goods",
    "SHAKTIPUMP.NS": "Capital Goods",
    "WABAG.NS": "Capital Goods",
    "APLAPOLLO.NS": "Steel",
    # Consumer
    "VBL.NS": "FMCG",
    "NESTLEIND.NS": "FMCG",
    "TRAVELFOOD.NS": "FMCG",
    # Real Estate / REIT
    "GODREJPROP.NS": "Real Estate",
    "MINDSPACE-RR.NS": "Real Estate",
    "EMBASSY.NS": "Real Estate",
    "BROOKFIELD.NS": "Real Estate",
    "NEXUS.NS": "Real Estate",
    # ETFs / commodity ETFs
    "GOLDBEES.NS": "Gold ETF",
    "SILVERBEES.NS": "Silver ETF",
    "MODEFENCE.NS": "Defence ETF",
    # Healthcare / Pharma
    "LAURUSLABS.NS": "Pharma",
    "GRANULES.NS": "Pharma",
    "NAVINFLUOR.NS": "Chemicals",
    "EMCURE.NS": "Pharma",
    "SAGILITY.NS": "IT",
    "HCLTECH.NS": "IT",
    # Chemicals / Materials
    "GRAVITA.NS": "Recycling",
    "GALAXYSURF.NS": "Chemicals",
    "APARINDS.NS": "Capital Goods",
    # Metals
    "NATIONALUM.NS": "Metals",
    # Technology / Electronics
    "DIXON.NS": "EMS",
    "NETWEB.NS": "IT",
    # Defence
    "MODEFENCE.NS": "Defence",
    # FMCG / Agri
    "CCL.NS": "FMCG",
    # Energy / Renewables
    "WAAREEENER.NS": "Renewable",
    # Misc
    "MANORAMA.NS": "FMCG",
    "SHAILY.NS": "Chemicals",
    "NH.NS": "Healthcare",
    "NEPHROPLUS.NS": "Healthcare",
    "ULTRACEMCO.NS": "Cement",
    "SAILIFE.NS": "Pharma",
    "LGEINDIA.NS": "Consumer Durables",
    "SETL.NS": "Capital Goods",
    "SCHAEFFLER.NS": "Auto Ancillary",
    "TVSMOTOR.NS": "Auto",
    "ENRIN.NS": "Energy",
    "RELIANCE.NS": "Energy",
    "QPOWER.NS": "Energy",
    "BLS.NS": "Travel",
    "ETERNAL.NS": "E-Commerce",
    "NYKAA.NS": "E-Commerce",
    "HINDALCO.NS": "Metals",
    "ITBEES.NS": "IT",
    "KARURVYSYA.NS": "BFSI",
    "SWSOLAR.NS": "Renewable",
    "TITAN.NS": "Gold",
    "VENUSPIPES.NS": "Capital Goods",
    "VIKRAN.NS": "Capital Goods",
    "CARTRADE.NS": "E-Commerce",
    "ASTERDM.NS": "Healthcare",
    "ASTRAMICRO.NS": "Defence",
    "BELRISE.NS": "Auto",
    "BSE.NS": "BFSI",
    "CMPDI.NS": "Mining",
    "KMEW.NS": "Shipping",
    "NEULANDLAB.NS": "Pharma",
    "RUBICON.NS": "Pharma",
    "SYRMA.NS": "EMS",
    "EMMVEE.NS": "Renewable",
    "RATEGAIN.NS": "Travel",
    "CPPLUS.NS": "Consumer Durables",
    "GROWW.NS": "BFSI",
    "INOXINDIA.NS": "Capital Goods",
    "RADICO.NS": "FMCG",
    "HDBFS.NS": "BFSI",
    "PRIVISCL.NS": "Chemicals",
    "SUNDRMFAST.NS": "Auto",
    "CGCL.NS": "BFSI",
    "JYOTICNC.NS": "Capital Goods",
    "TIINDIA.NS": "Auto",
    "JIOFIN.NS": "BFSI",
}


class TransactionsFormatError(ValueError):
    pass


def format_inr(x) -> str:
    """
    Indian-style digit grouping (e.g. 12,34,567 not 1,234,567) — the last
    three digits are grouped normally, then every two digits after that
    (lakhs, crores). Plain f"{x:,.0f}" formatting groups in international
    thousands and reads wrong for INR figures.
    """
    if x is None or (isinstance(x, float) and pd.isna(x)):
        return ""
    try:
        v = round(float(x))
    except (TypeError, ValueError):
        return ""

    sign = "-" if v < 0 else ""
    v = abs(v)
    s = str(int(v))

    if len(s) <= 3:
        result = s
    else:
        result = s[-3:]
        s = s[:-3]
        while len(s) > 2:
            result = s[-2:] + "," + result
            s = s[:-2]
        if s:
            result = s + "," + result

    return f"{sign}₹{result}"


def _normalize_columns(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    rename: dict[str, str] = {}
    for col in df.columns:
        rename[col] = str(col).strip().lower()
    df = df.rename(columns=rename)

    alias = {
        "client": "portfolio",
        "name": "portfolio",
        "dt": "date",
        "transaction_date": "date",
        "qty": "quantity",
        "units": "quantity",
        "rate": "price",
        "amount": "price",
        "symbol": "ticker",
    }
    df = df.rename(columns={k: v for k, v in alias.items() if k in df.columns})
    return df


def _normalize_columns_basic(df: pd.DataFrame) -> pd.DataFrame:
    """
    Lowercase/strip column names only — no alias remapping.

    Used for sheets (like Cashflows) whose columns must NOT go through the
    transaction-sheet alias map, since that map redirects "amount" -> "price"
    (meant for transaction sheets that call the trade value "Amount"), which
    would otherwise silently destroy the Cashflows "Amount" column.
    """
    df = df.copy()
    df = df.rename(columns={col: str(col).strip().lower() for col in df.columns})
    return df


def _find_sheet(sheets: dict[str, pd.DataFrame], *names: str) -> pd.DataFrame:
    """Case/whitespace-insensitive lookup of a worksheet by name."""
    lookup = {str(k).strip().lower(): k for k in sheets.keys()}
    for name in names:
        key = name.strip().lower()
        if key in lookup:
            return sheets[lookup[key]]
    raise TransactionsFormatError(
        f"Could not find a sheet named one of {list(names)}. "
        f"Found sheets: {list(sheets.keys())}"
    )


def load_transactions_excel(excel_bytes: bytes, password: Optional[str] = None) -> dict[str, pd.DataFrame]:
    """
    Load the full workbook (all sheets) as a dict {sheet_name: DataFrame}.

    The transactions master file has three sheets:
      - "Initial Holdings": positions each portfolio held before the tracked
        transaction history begins.
      - "Cashflows": deposits/withdrawals per portfolio.
      - "Transactions": the buy/sell trade log.
    """
    def _read(bio: io.BytesIO) -> dict[str, pd.DataFrame]:
        return pd.read_excel(bio, sheet_name=None)

    try:
        return _read(io.BytesIO(excel_bytes))
    except Exception:
        if not password:
            raise

    import msoffcrypto

    decrypted = io.BytesIO()
    office_file = msoffcrypto.OfficeFile(io.BytesIO(excel_bytes))
    office_file.load_key(password=password)
    office_file.decrypt(decrypted)
    decrypted.seek(0)
    return _read(decrypted)


def prepare_transactions(sheets: dict[str, pd.DataFrame]) -> tuple[pd.DataFrame, pd.DataFrame]:
    """
    Build the working transactions DataFrame and a separate cashflows DataFrame
    from the raw workbook sheets.

    "Initial Holdings" rows are converted into synthetic opening BUY transactions
    dated at each portfolio's StartDate, priced at AvgCostPrice. This lets FIFO
    cost-basis, daily-position, NAV, and attribution logic work correctly from
    each portfolio's true inception rather than only from its first *recorded*
    trade — without this, any position that predates the transaction log would
    be invisible to the rest of the pipeline (wrong holdings, wrong P&L, wrong
    NAV/returns).

    Returns:
        (transactions_df, cashflows_df)
    """
    holdings_raw = _find_sheet(sheets, "Initial Holdings", "InitialHoldings", "Holdings")
    txns_raw = _find_sheet(sheets, "Transactions")
    cash_raw = _find_sheet(sheets, "Cashflows", "Cash Flows", "CashFlows")

    # ---- Initial Holdings -> synthetic opening BUY transactions ----
    h = _normalize_columns(holdings_raw)
    missing_h = REQUIRED_HOLDINGS_COLUMNS - set(h.columns)
    if missing_h:
        raise TransactionsFormatError(
            f"'Initial Holdings' sheet missing columns: {sorted(missing_h)}. Found: {sorted(h.columns)}"
        )
    h_opening = pd.DataFrame(
        {
            "portfolio": h["portfolio"],
            "date": h["startdate"],
            "ticker": h["ticker"],
            "action": "BUY",
            "quantity": h["quantity"],
            "price": h["avgcostprice"],
            "is_opening": True,
            # Optional: carried through so build_daily_ledger can bootstrap Day-0
            # cash as (total opening value - equity value of these holdings) —
            # the same approach the reference PortfolioTWR.py uses — instead of
            # assuming zero starting cash when this column isn't present.
            "opening_market_value": h["marketvalue"] if "marketvalue" in h.columns else np.nan,
        }
    )

    # ---- Transactions ----
    t = _normalize_columns(txns_raw)
    missing_t = REQUIRED_COLUMNS - set(t.columns)
    if missing_t:
        raise TransactionsFormatError(
            f"'Transactions' sheet missing columns: {sorted(missing_t)}. Found: {sorted(t.columns)}"
        )
    t = t[list(REQUIRED_COLUMNS)].copy()
    t["is_opening"] = False
    t["opening_market_value"] = np.nan

    combined = pd.concat([h_opening, t], ignore_index=True, sort=False)

    combined["portfolio"] = combined["portfolio"].astype(str).str.strip()
    combined["ticker"] = combined["ticker"].astype(str).str.strip().str.upper()
    combined["action"] = combined["action"].astype(str).str.strip().str.upper()

    combined["date"] = pd.to_datetime(combined["date"], errors="coerce")
    combined = combined.dropna(subset=["date"])
    combined["date"] = combined["date"].dt.normalize()

    combined["quantity"] = pd.to_numeric(combined["quantity"], errors="coerce")
    combined["price"] = pd.to_numeric(combined["price"], errors="coerce")
    combined = combined.dropna(subset=["quantity", "price"])

    combined["signed_qty"] = np.where(combined["action"] == "BUY", combined["quantity"], -combined["quantity"])

    combined["ticker_yf"] = (
        combined["ticker"]
        .astype(str)
        .str.upper()
        .str.strip()
        .map(TICKER_MAP)
    )
    # Only append .NS to tickers that don't already carry an exchange suffix.
    # The previous version appended .NS unconditionally to anything not in
    # TICKER_MAP — if the sheet's own Ticker column already had ".NS" (or
    # ".BO") on some rows, that produced "RELIANCE.NS.NS": a symbol that
    # will NEVER resolve on Yahoo Finance no matter how many times it's
    # retried. This is very likely the actual, deterministic cause of the
    # "same ~40 tickers always fail" pattern — it's a malformed symbol, not
    # a network/rate-limit issue, which is exactly why retries never helped.
    _raw_ticker = combined["ticker"].astype(str).str.upper().str.strip()
    _already_suffixed = _raw_ticker.str.endswith(".NS") | _raw_ticker.str.endswith(".BO")
    _fallback = pd.Series(np.where(_already_suffixed, _raw_ticker, _raw_ticker + ".NS"), index=combined.index)
    combined["ticker_yf"] = combined["ticker_yf"].fillna(_fallback)

    combined = combined.sort_values(["portfolio", "date", "ticker_yf"]).reset_index(drop=True)

    # ---- Cashflows ----
    c = _normalize_columns_basic(cash_raw)
    missing_c = REQUIRED_CASHFLOW_COLUMNS - set(c.columns)
    if missing_c:
        raise TransactionsFormatError(
            f"'Cashflows' sheet missing columns: {sorted(missing_c)}. Found: {sorted(c.columns)}"
        )
    c = c[list(REQUIRED_CASHFLOW_COLUMNS)].copy()
    c["portfolio"] = c["portfolio"].astype(str).str.strip()
    c["date"] = pd.to_datetime(c["date"], errors="coerce")
    c = c.dropna(subset=["date"])
    c["date"] = c["date"].dt.normalize()
    c["amount"] = pd.to_numeric(c["amount"], errors="coerce")
    c = c.dropna(subset=["amount"])
    c = c.sort_values(["portfolio", "date"]).reset_index(drop=True)

    return combined, c


def cashflow_summary(cashflows_df: pd.DataFrame) -> pd.Series:
    """Total net cash contributed per portfolio (sum of Cashflows.Amount)."""
    if cashflows_df is None or cashflows_df.empty:
        return pd.Series(dtype=float)
    return cashflows_df.groupby("portfolio")["amount"].sum().sort_values(ascending=False)


def current_positions(df: pd.DataFrame) -> pd.DataFrame:
    """
    Return open positions per portfolio with qty and average buy price (FIFO).
    Only tickers with net qty > 0 (still held) are returned.
    """
    rows = []
    for portfolio in sorted(df["portfolio"].unique()):
        txns = df[df["portfolio"] == portfolio].sort_values("date")
        # FIFO lots: ticker -> [(qty, price), ...]
        lots: dict[str, list[tuple[float, float]]] = {}
        for _, row in txns.iterrows():
            ticker = str(row["ticker_yf"])
            qty    = float(row["quantity"])
            price  = float(row["price"])
            action = str(row["action"]).upper()
            if action == "BUY":
                lots.setdefault(ticker, []).append((qty, price))
            elif action == "SELL":
                remaining = qty
                old_lots  = lots.get(ticker, [])
                new_lots  = []
                for lq, lp in old_lots:
                    if remaining <= 0:
                        new_lots.append((lq, lp))
                    elif lq <= remaining:
                        remaining -= lq          # lot fully consumed
                    else:
                        new_lots.append((lq - remaining, lp))
                        remaining = 0
                lots[ticker] = new_lots

        for ticker, ticker_lots in lots.items():
            open_qty = sum(q for q, _ in ticker_lots)
            if open_qty <= 0:
                continue
            avg_price = sum(q * p for q, p in ticker_lots) / open_qty
            rows.append({
                "portfolio":   portfolio,
                "ticker_yf":   ticker,
                "qty":         round(open_qty, 4),
                "avg_buy_price": round(avg_price, 2),
            })

    if not rows:
        return pd.DataFrame(columns=["portfolio", "ticker_yf", "qty", "avg_buy_price"])

    out = pd.DataFrame(rows)
    return out.sort_values(["portfolio", "qty"], ascending=[True, False]).reset_index(drop=True)


def bse_ticker_from_ns(ticker_ns: str) -> str:
    base = ticker_ns.replace(".NS", "").replace(".BO", "")
    base = BSE_OVERRIDE.get(base, base)
    return f"{base}.BO"


def build_twr_returns(
    daily_positions: dict[str, pd.DataFrame],
    prices: pd.DataFrame,
) -> pd.DataFrame:
    """
    Compute daily Time-Weighted Returns (TWR) for each portfolio.

    For each day t:
        r_t = sum(holdings_{t-1} * price_t) / sum(holdings_{t-1} * price_{t-1}) - 1

    This uses YESTERDAY'S holdings with today's and yesterday's prices so that
    cash inflows/outflows (new buys/rebalancing sells) do NOT inflate the return.

    Days that are invalid for return measurement are set to NaN:
      - Days before the portfolio's first real position
      - Days where yesterday's NAV was below 2% of the portfolio's median NAV
        (these are rebalancing/entry days where the portfolio was temporarily
        fully-or-mostly-cash and the ratio produces nonsense)
    After filtering, returns are capped at ±25% per day — the hard limit for
    a diversified equity portfolio in a single session.
    """
    result = pd.DataFrame(index=prices.index)

    for portfolio, holdings in daily_positions.items():
        available = [t for t in holdings.columns if t in prices.columns]
        if not available:
            result[portfolio] = np.nan
            continue

        h = holdings[available].reindex(prices.index).ffill().fillna(0)
        p = prices[available].reindex(prices.index)

        prev_h = h.shift(1)
        nav_today     = (prev_h * p).sum(axis=1)
        nav_yesterday = (prev_h * p.shift(1)).sum(axis=1)

        # NaN out days where prior-day NAV was zero or below 2% of median NAV.
        # This handles full rebalancing days (sell-all → rebuy) where the prior
        # holding is near zero and produces a meaningless giant return.
        median_nav = nav_yesterday[nav_yesterday > 0].median()
        min_nav_threshold = median_nav * 0.02 if pd.notna(median_nav) and median_nav > 0 else 0
        valid_base = nav_yesterday.where(nav_yesterday >= min_nav_threshold)

        twr = (nav_today / valid_base - 1)
        twr = twr.replace([np.inf, -np.inf], np.nan)

        # NaN out everything before the first day the portfolio had real holdings
        first_valid = h[(h > 0).any(axis=1)].index
        if len(first_valid):
            twr.loc[twr.index < first_valid[0]] = np.nan

        # Hard cap: no equity portfolio moves more than ±25% in a single day
        twr = twr.clip(lower=-0.25, upper=0.25)

        result[portfolio] = twr

    return result.dropna(how="all")


def build_daily_ledger(
    df: pd.DataFrame, cashflows_df: pd.DataFrame, prices: pd.DataFrame
) -> tuple[dict[str, pd.DataFrame], dict[str, pd.DataFrame], pd.DataFrame]:
    """
    Full daily ledger per portfolio — Equity Value, Cash, Portfolio Value,
    and External Flow — walked forward day by day. This is a direct port of
    the reference PortfolioTWR.py's `build_daily_ledger_existing` (for
    portfolios with Initial Holdings) and `build_daily_ledger` (fresh-start,
    no Initial Holdings) — including their specific Day-0 handling, not
    just the general idea of "ledger with cash".

    This is now also the ONE AND ONLY place per-day share quantities are
    tracked for the whole app — attribution, risk contribution, sector
    exposure, market-cap allocation, and the holdings table all read from
    the same `daily_positions` this function returns, instead of each
    maintaining (and potentially disagreeing with) its own separate
    position tracker. There used to be three independent implementations
    of "what do we hold on date X" in this codebase (this ledger's own
    walk, a separate build_daily_positions(), and current_positions()'s
    FIFO lots) — they were never guaranteed to agree with each other, and
    the cross-check between them produced nothing but unexplained
    "disagreement" errors. current_positions()/FIFO lots still exist
    separately because they answer a genuinely different question (cost
    basis for P&L, which requires lot-level tracking that a quantity-only
    ledger can't provide) — but there's no longer a second, parallel
    quantity tracker for the same thing.

    For a portfolio WITH Initial Holdings:
      - Day 0 (StartDate) is recorded as a baseline row using the sheet's
        own numbers verbatim: Equity Value = Portfolio Value = sum of
        Initial Holdings' MarketValue, Cash = 0. This guarantees the very
        first data point matches the sheet exactly, with no
        computation-derived rounding/gap.
      - Internally (not shown in that row), cash is bootstrapped as
        (MarketValue total) - (those quantities valued at StartDate
        price), and any Cashflow/Transaction dated exactly on StartDate is
        applied to that internal cash/holdings state right after the
        baseline row.
      - Because that Day-0 cash impact isn't visible in the Day-0 row, it's
        carried forward and added to the FIRST real row's External Flow —
        otherwise compute_ledger_twr would misread the jump from the raw
        baseline to the adjusted value as investment return.
      - From the first real row onward, this behaves like the fresh-start
        case below.

    For a portfolio WITHOUT Initial Holdings (fresh start):
      - Cash starts at 0.0 and is built up purely by Cashflows deposits and
        by BUY/SELL trade cash impact, from the portfolio's first activity
        date.

    Every SELL is clamped so a ticker's quantity never goes negative — a
    purely defensive safety net at this point, since `df` is expected to
    already have gone through `auto_backfill_oversold_positions`, which
    should make a genuine oversold event impossible. Any clamp that DOES
    fire is returned in `residual_oversold` so it's visible rather than
    silent, since it would mean the backfill step itself has a gap.

    Returns (ledgers, daily_positions, residual_oversold):
      - ledgers: {portfolio: DataFrame(index=date, columns=[equity_value,
        cash, portfolio_value, external_flow])}
      - daily_positions: {portfolio: DataFrame(index=date, columns=ticker_yf)}
      - residual_oversold: DataFrame of any SELL that still had to be
        clamped despite the upstream backfill (expected to be empty).
    """
    ledgers: dict[str, pd.DataFrame] = {}
    daily_positions: dict[str, pd.DataFrame] = {}
    residual_oversold_rows: list[dict] = []
    has_is_opening = "is_opening" in df.columns
    has_opening_mv = "opening_market_value" in df.columns

    for portfolio in sorted(df["portfolio"].unique().tolist()):
        txns = df[df["portfolio"] == portfolio].copy()
        cfs = (
            cashflows_df[cashflows_df["portfolio"] == portfolio].copy()
            if cashflows_df is not None and not cashflows_df.empty
            else pd.DataFrame(columns=["date", "amount"])
        )

        opening_rows = (
            txns[txns["is_opening"] == True] if has_is_opening else pd.DataFrame()  # noqa: E712
        )
        opening_mv_sum = (
            pd.to_numeric(opening_rows["opening_market_value"], errors="coerce").sum()
            if has_opening_mv and not opening_rows.empty
            else np.nan
        )
        has_existing_holdings = not opening_rows.empty and pd.notna(opening_mv_sum) and opening_mv_sum > 0

        start_candidates = []
        if has_existing_holdings:
            start_candidates.append(opening_rows["date"].iloc[0])
        if not txns.empty:
            start_candidates.append(txns["date"].min())
        if not cfs.empty:
            start_candidates.append(cfs["date"].min())
        if not start_candidates:
            continue
        start = min(start_candidates)
        end = prices.index.max()

        activity_dates = pd.DatetimeIndex(
            pd.concat([txns["date"], cfs["date"]], ignore_index=True).dropna().unique()
        )
        market_dates = pd.DatetimeIndex(prices.index).normalize()
        dates = market_dates.union(activity_dates).union(pd.DatetimeIndex([start])).sort_values()
        dates = dates[(dates >= start) & (dates <= end)]
        if len(dates) == 0:
            continue

        # bfill (in addition to ffill) only for the existing-holdings branch,
        # matching the reference exactly — it needs a valid StartDate price
        # even for a ticker that only starts trading slightly later.
        p = prices.reindex(dates).ffill().bfill() if has_existing_holdings else prices.reindex(dates).ffill()

        holdings: dict[str, float] = {t: 0.0 for t in txns["ticker_yf"].unique()}
        txn_by_date = {d: g for d, g in txns.groupby("date")}
        cf_by_date = cfs.groupby("date")["amount"].sum().to_dict() if not cfs.empty else {}

        def _apply_sell(ticker: str, qty: float, txn_date) -> None:
            available = holdings.get(ticker, 0.0)
            actual = min(qty, available)
            oversold = qty - actual
            holdings[ticker] = available - actual
            if oversold > 1e-9:
                residual_oversold_rows.append(
                    {"portfolio": portfolio, "ticker_yf": ticker, "date": txn_date, "oversold_qty": round(oversold, 4)}
                )

        rows: list[dict] = []
        position_rows: list[dict] = []

        if has_existing_holdings:
            opening_qty = opening_rows.groupby("ticker_yf")["quantity"].sum()
            for tkr, qty in opening_qty.items():
                holdings[tkr] = holdings.get(tkr, 0.0) + float(qty)
            start_prices = p.iloc[0] if not p.empty else pd.Series(dtype=float)
            opening_equity = sum(
                float(qty) * float(start_prices[tkr])
                for tkr, qty in opening_qty.items()
                if tkr in start_prices.index and pd.notna(start_prices.get(tkr))
            )
            opening_total = float(opening_mv_sum)
            cash = opening_total - opening_equity

            # Day-0 baseline row: the sheet's own numbers, verbatim.
            rows.append(
                {"date": start, "equity_value": opening_total, "cash": 0.0,
                 "portfolio_value": opening_total, "external_flow": 0.0}
            )
            position_rows.append({"date": start, **holdings})

            # Apply any StartDate-dated activity to the internal state AFTER
            # recording the baseline, so it affects every subsequent day.
            cash += float(cf_by_date.get(start, 0.0))
            for _, r in txn_by_date.get(start, pd.DataFrame()).iterrows():
                if bool(r.get("is_opening", False)):
                    continue  # already folded into the opening_qty bootstrap above
                ticker = r["ticker_yf"]
                qty = float(r["quantity"])
                value = qty * float(r["price"])
                if str(r["action"]).upper() == "BUY":
                    holdings[ticker] = holdings.get(ticker, 0.0) + qty
                    cash -= value
                else:
                    _apply_sell(ticker, qty, start)
                    cash += value

            remaining_dates = dates[dates > start]
        else:
            cash = 0.0
            remaining_dates = dates

        for d in remaining_dates:
            if d in cf_by_date:
                cash += float(cf_by_date[d])

            if d in txn_by_date:
                for _, r in txn_by_date[d].iterrows():
                    if bool(r.get("is_opening", False)):
                        continue  # opening rows only ever apply on `start`, handled above
                    ticker = r["ticker_yf"]
                    qty = float(r["quantity"])
                    value = qty * float(r["price"])
                    if str(r["action"]).upper() == "BUY":
                        holdings[ticker] = holdings.get(ticker, 0.0) + qty
                        cash -= value
                    else:
                        _apply_sell(ticker, qty, d)
                        cash += value

            equity = 0.0
            for ticker, qty in holdings.items():
                if abs(qty) < 1e-9:
                    continue
                if ticker in p.columns:
                    px = p.loc[d, ticker]
                    if pd.notna(px):
                        equity += qty * px

            # The Day-0 cash impact of a StartDate-dated flow isn't visible in
            # the baseline row, so it's carried forward onto the first real
            # row's External Flow — otherwise compute_ledger_twr would read
            # the jump from the raw baseline as investment return.
            start_flow = 0.0
            if has_existing_holdings and len(rows) == 1:
                start_flow = float(cf_by_date.get(start, 0.0))

            rows.append(
                {
                    "date": d,
                    "equity_value": equity,
                    "cash": cash,
                    "portfolio_value": equity + cash,
                    "external_flow": float(cf_by_date.get(d, 0.0)) + start_flow,
                }
            )
            position_rows.append({"date": d, **holdings})

        ledgers[str(portfolio)] = pd.DataFrame(rows).set_index("date")
        pos_df = pd.DataFrame(position_rows).set_index("date").fillna(0.0)
        # Restrict the returned position history to the actual price
        # calendar (the baseline row above may predate it), same convention
        # build_daily_positions used previously.
        daily_positions[str(portfolio)] = pos_df.reindex(prices.index).ffill().fillna(0.0)

    residual_oversold = pd.DataFrame(
        residual_oversold_rows, columns=["portfolio", "ticker_yf", "date", "oversold_qty"]
    )
    return ledgers, daily_positions, residual_oversold


def compute_ledger_twr(ledger: pd.DataFrame) -> pd.Series:
    """
    Daily TWR return from a full ledger (equity + cash combined) — a direct
    port of the reference PortfolioTWR.py's `compute_twr`:

        r_t = (portfolio_value_t - external_flow_t) / portfolio_value_{t-1} - 1

    Subtracting the day's external flow (a Cashflows deposit/withdrawal)
    before dividing means a deposit or withdrawal never shows up as a gain
    or loss.

    Matches the reference exactly: a day's return defaults to 0.0 (not
    NaN) whenever the prior day's portfolio value was at or below zero —
    there's no ±% clip and no special-casing around negative cash. This is
    deliberately simpler than (and supersedes) an earlier version of this
    function that added those guards; the guards were this app's own
    invention and a source of divergence from the reference, not a
    correction to it.
    """
    if ledger is None or ledger.empty:
        return pd.Series(dtype=float)

    d = ledger.sort_index()
    pv = pd.to_numeric(d["portfolio_value"], errors="coerce")
    flow = pd.to_numeric(d["external_flow"], errors="coerce").fillna(0.0)
    prev_pv = pv.shift(1)

    ret = pd.Series(0.0, index=d.index)
    valid = prev_pv > 0
    ret.loc[valid] = (pv.loc[valid] - flow.loc[valid]) / prev_pv.loc[valid] - 1
    ret = ret.replace([np.inf, -np.inf], 0.0)
    return ret


def annualise_twr_simple(twr: float, days: int) -> float:
    """Direct port of the reference's `annualise_twr` — annualizes a single
    cumulative TWR figure over the actual number of calendar days elapsed."""
    if days is None or days <= 0 or pd.isna(twr):
        return np.nan
    return float((1 + twr) ** (365 / days) - 1)


def trailing_returns_table(pr: pd.Series, benchmark_price_series: dict[str, pd.Series]) -> pd.DataFrame:
    """
    Standard trailing-period returns table — Portfolio TWR vs a benchmark's
    own price return, for 1M / 3M / 6M / 1Y / Since Inception, plus an
    Alpha row (Portfolio minus benchmark, period by period).

    Returns a DataFrame with rows = Portfolio, <benchmark name>, Alpha and
    columns = 1M/3M/6M/1Y/Since Inception (periods as columns) — only the
    first benchmark passed in is used, since this table is meant to show
    one primary comparison, not several.

    Portfolio uses its own TWR growth-of-1 curve (built from daily TWR
    returns, so cash flows never distort it). The benchmark uses its own
    plain price return over the same window — not the cashflow-replicated
    simulation used for the Alpha figures elsewhere, since a trailing
    comparison table is meant to show "how did the index itself do over
    this window", which is what every factsheet means by it.

    A period with insufficient history (the portfolio or the benchmark
    didn't exist that far back) is left as NaN — rendered blank by the
    caller — rather than computed on a truncated window or raising.
    """
    pr = pd.to_numeric(pr, errors="coerce").dropna()
    if pr.empty:
        return pd.DataFrame()

    port_growth = (1 + pr).cumprod()
    end_date = port_growth.index.max()

    series_map: dict[str, pd.Series] = {"Portfolio": port_growth}
    bench_name = None
    for name, s in benchmark_price_series.items():
        s = pd.to_numeric(s, errors="coerce").dropna()
        if s.empty:
            continue
        series_map[name] = s / s.iloc[0]
        bench_name = name
        break  # only the first (primary) benchmark — one comparison, not several

    period_offsets = [
        ("1M", pd.DateOffset(months=1)),
        ("3M", pd.DateOffset(months=3)),
        ("6M", pd.DateOffset(months=6)),
        ("1Y", pd.DateOffset(years=1)),
    ]
    period_labels = [lbl for lbl, _ in period_offsets] + ["Since Inception"]

    rows: dict[str, dict] = {name: {} for name in series_map}
    for label, offset in period_offsets:
        cutoff = end_date - offset
        for col, growth in series_map.items():
            if growth.empty or growth.index.min() > cutoff:
                rows[col][label] = np.nan  # didn't exist that far back — leave blank
                continue
            start_val = growth.asof(cutoff)
            end_val = growth.iloc[-1]
            rows[col][label] = np.nan if pd.isna(start_val) or start_val == 0 else float(end_val / start_val - 1)

    for col, growth in series_map.items():
        rows[col]["Since Inception"] = np.nan if growth.empty else float(growth.iloc[-1] / growth.iloc[0] - 1)

    table = pd.DataFrame(rows).T  # rows = series name, columns = period label
    table = table.reindex(columns=period_labels)

    if bench_name is not None:
        table.loc["Alpha"] = table.loc["Portfolio"] - table.loc[bench_name]

    return table


def benchmark_replication(cashflows_df: pd.DataFrame, series: pd.Series) -> pd.DataFrame:
    """
    Direct port of the reference's `benchmark_replication`: simulates
    investing the SAME external cashflows into the benchmark instead, by
    buying benchmark "units" with each flow at that day's benchmark price.
    This is what makes the resulting Alpha (portfolio TWR - benchmark TWR)
    a fair like-for-like comparison — same cashflow timing on both sides —
    rather than comparing against a benchmark that was simply bought once
    on day one.
    """
    if cashflows_df is None or cashflows_df.empty:
        flows: dict = {}
    else:
        flows = cashflows_df.groupby("date")["amount"].sum().to_dict()

    units = 0.0
    rows = []
    for d in series.index:
        px = series.loc[d]
        flow = flows.get(d, 0.0)
        if flow != 0 and pd.notna(px) and px > 0:
            units += flow / px
        value = units * px if pd.notna(px) else 0.0
        rows.append({"date": d, "equity_value": value, "cash": 0.0, "portfolio_value": value, "external_flow": flow})
    return pd.DataFrame(rows).set_index("date")


def benchmark_existing(cashflows_df: pd.DataFrame, series: pd.Series, opening: float, start) -> pd.DataFrame:
    """Direct port of the reference's `benchmark_existing`: same idea as
    benchmark_replication, but seeded with `opening` value's worth of units
    on `start` (matching an existing-holdings portfolio's own Day-0
    baseline) instead of starting from zero."""
    series = series[series.index >= start].ffill().bfill()
    if series.empty or pd.isna(series.iloc[0]) or series.iloc[0] <= 0:
        return pd.DataFrame(columns=["equity_value", "cash", "portfolio_value", "external_flow"])
    flows = (
        cashflows_df[cashflows_df["date"] >= start].groupby("date")["amount"].sum().to_dict()
        if cashflows_df is not None and not cashflows_df.empty
        else {}
    )
    units = opening / float(series.iloc[0])
    start_flow = float(flows.get(start, 0.0))
    rows = []
    for i, (d, px) in enumerate(series.items()):
        if i == 0:
            flow = 0.0
            value = units * float(px)
        else:
            flow = float(flows.get(d, 0.0))
            if i == 1 and start_flow:
                flow += start_flow
            if flow:
                units += flow / float(px)
            value = units * float(px)
        rows.append({"date": d, "equity_value": value, "cash": 0.0, "portfolio_value": value, "external_flow": flow})
    return pd.DataFrame(rows).set_index("date")


def annualized_return_from_twr(twr_returns: pd.Series) -> float:
    """
    Annualized return from a series of daily TWR returns.

    Uses actual calendar-day scaling:
        (prod(1+r_i))^(365 / calendar_days) - 1

    Falls back to trading-day scaling (252) if the index has no DatetimeIndex.
    """
    r = pd.to_numeric(twr_returns, errors="coerce").dropna()
    r = r.replace([np.inf, -np.inf], np.nan).dropna()
    n = len(r)
    if n < 2:
        return np.nan
    cumulative = float((1 + r).prod())
    if cumulative <= 0:
        return np.nan

    # Prefer calendar-day scaling when index carries dates
    if isinstance(r.index, pd.DatetimeIndex) and len(r.index) >= 2:
        cal_days = (r.index[-1] - r.index[0]).days
        if cal_days >= 2:
            return float(cumulative ** (365.0 / cal_days) - 1)

    # Fallback: trading-day scaling
    return float(cumulative ** (252.0 / n) - 1)


def cumulative_growth_from_twr(twr_returns: pd.Series) -> pd.Series:
    """
    Growth-of-₹1 index built purely from daily TWR returns (no NAV/dollar value
    involved). Used for max-drawdown and the PDF growth chart. Because it's
    built from TWR (which already excludes cash-flow effects), a drawdown
    computed off this curve reflects actual investment performance — unlike a
    drawdown computed off raw NAV, which would show a fake "loss" whenever the
    client simply withdrew cash.
    """
    r = pd.to_numeric(twr_returns, errors="coerce").dropna()
    if r.empty:
        return pd.Series(dtype=float)
    return (1 + r).cumprod()


def annualized_volatility(returns: pd.Series) -> float:
    r = pd.to_numeric(returns, errors="coerce").dropna()
    if len(r) < 2:
        return np.nan
    return float(r.std(ddof=1) * np.sqrt(252))


def _period_days(r: pd.Series) -> int:
    """Calendar days spanned by a returns series, for scaling a rate to match."""
    if isinstance(r.index, pd.DatetimeIndex) and len(r.index) >= 2:
        return max(int((r.index[-1] - r.index[0]).days), 1)
    return max(len(r) - 1, 1)


def _period_rf(rf_annual: float, days: int) -> float:
    """Risk-free rate compounded down to the same span as the return period,
    rather than mixing an annual rate with a period return."""
    return float((1 + rf_annual) ** (days / 365) - 1)


def sortino_ratio(returns: pd.Series, rf_annual: float = 0.0) -> float:
    """
    TWR-period Sortino — the period's own cumulative TWR against the
    period's own downside deviation, not an annualized projection of
    either. A short, choppy reporting period shouldn't be stretched into
    an annualized number that implies a full year of that behavior.
    """
    r = pd.to_numeric(returns, errors="coerce").dropna()
    r = r.replace([np.inf, -np.inf], np.nan).dropna()
    n = len(r)
    if n < 2:
        return np.nan
    rf_daily = (1 + rf_annual) ** (1 / 252) - 1
    downside = r[r < rf_daily] - rf_daily
    if len(downside) < 2:
        return np.nan
    downside_dev = float(downside.std(ddof=1) * np.sqrt(n))
    if downside_dev <= 0:
        return np.nan
    period_twr = float((1 + r).prod() - 1)
    period_rf = _period_rf(rf_annual, _period_days(r))
    return float((period_twr - period_rf) / downside_dev)


def beta_to_benchmark(returns: pd.Series, benchmark_returns: pd.Series) -> float:
    r = pd.to_numeric(returns, errors="coerce")
    b = pd.to_numeric(benchmark_returns, errors="coerce")
    df = pd.concat([r.rename("p"), b.rename("b")], axis=1).dropna()
    if len(df) < 3:
        return np.nan
    var_b = float(df["b"].var(ddof=1))
    # A benchmark's daily variance corresponding to well under 1% annualized
    # volatility is a sign of stale/near-constant price data (e.g. a thinly
    # traded mutual-fund NAV proxy with long flat stretches from
    # forward-filling), not a genuinely low-risk benchmark. Dividing by a
    # near-zero variance blows Beta up to nonsensical magnitudes (50+),
    # which then cascades into equally nonsensical Jensen Alpha/Treynor —
    # better to report "not available" than a number with no real meaning.
    min_var_b = (0.01 / np.sqrt(252)) ** 2
    if var_b <= min_var_b:
        return np.nan
    cov = float(df["p"].cov(df["b"]))
    return float(cov / var_b)


def jensen_alpha(
    returns: pd.Series,
    benchmark_returns: pd.Series,
    rf_annual: float = 0.0,
) -> float:
    """
    TWR-period Jensen's Alpha — the period's own cumulative TWR against
    what CAPM would have predicted for that SAME period (given the
    portfolio's Beta and the benchmark's own period return), not an
    annualized projection of either side.
    """
    r = pd.to_numeric(returns, errors="coerce").dropna()
    b = pd.to_numeric(benchmark_returns, errors="coerce").dropna()
    if len(r) < 2 or len(b) < 2:
        return np.nan
    beta = beta_to_benchmark(r, b)
    if pd.isna(beta):
        return np.nan
    period_p = float((1 + r).prod() - 1)
    period_b = float((1 + b).prod() - 1)
    period_rf = _period_rf(rf_annual, _period_days(r))
    return float(period_p - (period_rf + beta * (period_b - period_rf)))


def treynor_ratio(
    returns: pd.Series,
    benchmark_returns: pd.Series,
    rf_annual: float = 0.0,
) -> float:
    """TWR-period Treynor — period return per unit of Beta, not annualized."""
    r = pd.to_numeric(returns, errors="coerce").dropna()
    b = pd.to_numeric(benchmark_returns, errors="coerce").dropna()
    if len(r) < 2 or len(b) < 2:
        return np.nan
    beta = beta_to_benchmark(r, b)
    if pd.isna(beta) or beta == 0:
        return np.nan
    period_p = float((1 + r).prod() - 1)
    period_rf = _period_rf(rf_annual, _period_days(r))
    return float((period_p - period_rf) / beta)


def cvar_99(returns: pd.Series) -> float:
    """
    99% Conditional Value at Risk (Expected Shortfall).
    Average of all returns below the 1st percentile, sign-flipped to positive.
    More informative than VaR — measures the expected loss *given* we're in the tail.
    """
    r = pd.to_numeric(returns, errors="coerce").dropna()
    if len(r) < 10:
        return np.nan
    cutoff = float(r.quantile(0.01))
    tail = r[r <= cutoff]
    if tail.empty:
        return np.nan
    return float(-tail.mean())


def portfolio_metrics(
    twr_returns: pd.Series,
    benchmarks: dict[str, pd.Series] | None = None,
    rf_annual: float = 0.0,
) -> pd.Series:
    """
    Compute portfolio metrics purely from daily Time-Weighted Returns.

    No NAV/dollar-value series is used anywhere here — annualized return,
    volatility, Sharpe, Sortino, CVaR, and max drawdown are all derived from
    `twr_returns`. This keeps every metric immune to cash inflows/outflows,
    since TWR already excludes them.

    `benchmarks` maps a display name to that benchmark's price series (e.g.
    {"Nifty 50": ..., "Nifty 500": ...}) — Beta, Jensen Alpha, and Treynor are
    computed against each one and returned as "Beta (<name>)", etc.
    """
    base_empty = {
        "Period Return (TWR)": np.nan,
        "Period Volatility (TWR)": np.nan,
        "Sharpe": np.nan,
        "Sortino": np.nan,
        "CVaR 99% (Daily)": np.nan,
        "Max Drawdown": np.nan,
    }
    benchmark_keys = {}
    for name in (benchmarks or {}).keys():
        benchmark_keys[f"Beta ({name})"] = np.nan
        benchmark_keys[f"Jensen Alpha ({name})"] = np.nan
        benchmark_keys[f"Treynor ({name})"] = np.nan

    rets = pd.to_numeric(twr_returns, errors="coerce").replace([np.inf, -np.inf], np.nan).dropna()
    if len(rets) < 2:
        return pd.Series({**base_empty, **benchmark_keys})

    # Period Return/Volatility are the ACTUAL TWR performance over this
    # specific reporting period — not a projected annual figure. Consistent
    # with Sharpe/Sortino/Alpha/Treynor below, which all use the same
    # period basis rather than mixing annualized and period figures.
    period_twr = float((1 + rets).prod() - 1)
    period_vol = float(rets.std(ddof=1) * np.sqrt(len(rets))) if len(rets) >= 2 else np.nan
    period_rf = _period_rf(rf_annual, _period_days(rets))
    sharpe = (period_twr - period_rf) / period_vol if period_vol and period_vol > 0 else np.nan

    # Max drawdown: from the TWR growth-of-₹1 curve, not raw NAV — a cash
    # withdrawal should never show up as a "drawdown".
    growth = cumulative_growth_from_twr(rets)
    if len(growth) >= 2:
        dd_series = drawdown_series(growth)
        max_dd = float(dd_series.min()) if not dd_series.empty else np.nan
    else:
        max_dd = np.nan

    out = {
        "Period Return (TWR)": period_twr,
        "Period Volatility (TWR)": period_vol,
        "Sharpe": sharpe,
        "Sortino": sortino_ratio(rets, rf_annual=rf_annual),
        "CVaR 99% (Daily)": cvar_99(rets),
        "Max Drawdown": max_dd,
    }
    out.update(benchmark_keys)

    for name, bench_prices in (benchmarks or {}).items():
        bprices = pd.to_numeric(bench_prices, errors="coerce").dropna()
        if len(bprices) >= 2:
            b_rets = bprices.pct_change().replace([np.inf, -np.inf], np.nan).dropna()
            out[f"Beta ({name})"] = beta_to_benchmark(rets, b_rets)
            out[f"Jensen Alpha ({name})"] = jensen_alpha(rets, b_rets, rf_annual=rf_annual)
            out[f"Treynor ({name})"] = treynor_ratio(rets, b_rets, rf_annual=rf_annual)

    return pd.Series(out)


@dataclass(frozen=True)
class DownloadResult:
    prices: pd.DataFrame
    successful_mappings: list[str]
    failed_tickers: list[str]
    failure_reasons: dict[str, str]


def download_prices(
    yf_module,
    tickers_ns: Iterable[str],
    start: datetime | date,
    include_benchmark: bool = True,
    benchmark_tickers: Iterable[str] | None = None,
    max_retries: int = 3,
    retry_delay: float = 1.0,
) -> DownloadResult:
    """
    Two independent yfinance code paths are tried per candidate symbol —
    `.download()` and `Ticker(...).history()` — before giving up on it,
    since the two go through different internal request logic and a
    failure specific to one (a schema change, a session/cookie issue) often
    doesn't affect the other. A ticker only lands in `failed` after both
    paths fail across every retry, for both the .NS and BSE candidate.

    Candidate order (.NS first, BSE fallback second) and the explicit `end`
    date with a 2-day buffer match the reference PortfolioTWR.py exactly —
    both are deliberate, not incidental.

    The actual exception (or "empty response") from the LAST attempt is
    kept per failed ticker in `failure_reasons` — a bare `except: pass`
    here means a systemic problem (a broken yfinance version, a
    blocked/rate-limited session) looks identical to a genuinely bad
    ticker.
    """
    import time

    all_prices: list[pd.Series] = []
    successful: list[str] = []
    failed: list[str] = []
    failure_reasons: dict[str, str] = {}
    end = pd.Timestamp.today().normalize() + pd.Timedelta(days=2)

    def _extract_close(df: pd.DataFrame, symbol: str):
        if df is None or getattr(df, "empty", True):
            return None
        if "Close" in df.columns:
            close = df["Close"]
        elif isinstance(df.columns, pd.MultiIndex) and ("Close", symbol) in df.columns:
            close = df[("Close", symbol)]
        else:
            return None
        if isinstance(close, pd.DataFrame):
            close = close.squeeze()
        if np.isscalar(close):
            return None
        s = pd.Series(close)
        return _tz_naive(s)

    def _tz_naive(s: Optional[pd.Series]) -> Optional[pd.Series]:
        """
        Strip tz info from a price series' index, keeping the wall-clock
        date/time as-is. `.download()` and `Ticker(...).history()` can
        return differently tz-localized indices (e.g. one tz-naive, one
        tagged Asia/Kolkata) depending on which succeeded for a given
        ticker — pd.concat() hard-fails ("Cannot join tz-naive with
        tz-aware DatetimeIndex") the moment two such series meet, so every
        series is normalized right where it's produced rather than trying
        to catch every combination at the concat call site.
        """
        if s is None:
            return None
        idx = s.index
        if isinstance(idx, pd.DatetimeIndex) and idx.tz is not None:
            s = s.copy()
            s.index = idx.tz_localize(None)
        return s

    def _download_with_retries(sym: str) -> tuple[Optional[pd.Series], Optional[str]]:
        last_err = "empty response"
        for attempt in range(max_retries):
            # Path 1: .download()
            try:
                temp = yf_module.download(sym, start=start, end=end, auto_adjust=True, progress=False)
                s = _extract_close(temp, sym)
                if s is not None and not s.dropna().empty:
                    return s, None
                last_err = "download() returned no rows"
            except Exception as e:
                last_err = f"download() raised {type(e).__name__}: {e}"

            # Path 2: Ticker(...).history() — a different internal code path,
            # so a failure specific to .download() often doesn't affect this.
            try:
                hist = yf_module.Ticker(sym).history(start=start, end=end, auto_adjust=True)
                if hist is not None and not hist.empty and "Close" in hist.columns:
                    s = _tz_naive(hist["Close"].dropna())
                    if s is not None and not s.empty:
                        return s, None
                last_err = "history() returned no rows"
            except Exception as e:
                last_err = f"history() raised {type(e).__name__}: {e}"

            if attempt < max_retries - 1:
                time.sleep(retry_delay)
        return None, last_err

    for i, ticker in enumerate(tickers_ns):
        candidates: list[str] = []
        # Reference-app order: .NS FIRST, BSE second. The BSE-first order
        # this app used before wastes the first attempt on a symbol that,
        # for a lot of names (ETFs especially — GOLDBEES, SILVERBEES, etc.),
        # doesn't reliably have BSE data at all — before ever reaching the
        # .NS symbol that actually works. Matching the reference's order is
        # a real behavior change, not cosmetic.
        candidates.append(ticker)
        candidates.append(NS_TO_BSE_TICKER.get(ticker, bse_ticker_from_ns(ticker)))

        got = None
        used = None
        last_reason = None
        for sym in candidates:
            s, err = _download_with_retries(sym)
            if s is not None:
                got = s
                used = sym
                break
            last_reason = f"{sym}: {err}"

        if i > 0:
            time.sleep(0.15)  # pacing delay between tickers

        if got is None or used is None:
            failed.append(ticker)
            if last_reason:
                failure_reasons[ticker] = last_reason
            continue

        got.name = ticker
        all_prices.append(got)
        successful.append(f"{ticker} -> {used}")

    prices = pd.concat(all_prices, axis=1) if all_prices else pd.DataFrame()

    if include_benchmark:
        benches = list(benchmark_tickers) if benchmark_tickers is not None else [BENCHMARK_TICKER]
        for bt in benches:
            bench_close, bench_err = _download_with_retries(bt)
            if bench_close is not None:
                prices[bt] = bench_close
            elif bench_err:
                failure_reasons[bt] = bench_err

    return DownloadResult(
        prices=prices, successful_mappings=successful, failed_tickers=failed, failure_reasons=failure_reasons
    )



def fetch_market_caps(yf_module, tickers: Iterable[str]) -> pd.Series:
    """
    Best-effort market cap (INR) per ticker via yfinance.

    ETFs are deliberately not assigned a company market cap. REITs use Yahoo's
    marketCap when available, with sharesOutstanding × current price as a
    fallback because marketCap is occasionally missing for REIT listings.
    """
    caps: dict[str, float] = {}

    def _try_one(symbol: str) -> Optional[float]:
        try:
            base = symbol.upper().replace(".BO", ".NS")
            if ASSET_CLASS_MAP.get(base) == "ETF":
                return None

            tk = yf_module.Ticker(symbol)
            mc = None
            shares = None
            price = None

            fast_info = getattr(tk, "fast_info", None)
            if fast_info is not None:
                try:
                    mc = fast_info.get("market_cap") if hasattr(fast_info, "get") else None
                except Exception:
                    mc = None
                if mc is None:
                    try:
                        mc = fast_info["market_cap"]
                    except Exception:
                        mc = None
                try:
                    price = fast_info.get("last_price") if hasattr(fast_info, "get") else None
                except Exception:
                    price = None

            try:
                info = tk.info
                if isinstance(info, dict):
                    if mc is None:
                        mc = info.get("marketCap")
                    shares = info.get("sharesOutstanding")
                    if price is None:
                        price = info.get("currentPrice") or info.get("regularMarketPrice")
            except Exception:
                pass

            # REIT fallback: market cap = units/shares outstanding × current price.
            if mc is None and shares is not None and price is not None:
                mc = float(shares) * float(price)

            if mc is not None and pd.notna(mc):
                return float(mc)
        except Exception:
            pass
        return None

    for ticker in tickers:
        mc = _try_one(ticker)
        if mc is None:
            alt = (
                NS_TO_BSE_TICKER.get(ticker, bse_ticker_from_ns(ticker))
                if ticker.endswith(".NS")
                else ticker.replace(".BO", ".NS")
            )
            if alt != ticker:
                mc = _try_one(alt)
        if mc is not None:
            caps[ticker] = mc

    return pd.Series(caps, dtype="float64")
def auto_backfill_oversold_positions(df: pd.DataFrame) -> tuple[pd.DataFrame, pd.DataFrame]:
    """
    For any (portfolio, ticker) where a SELL exceeds the quantity on record
    up to that point — no matching BUY in either Initial Holdings or prior
    Transactions — insert a synthetic BUY sized to cover exactly that
    shortfall, dated one day before the SELL that revealed it, priced at
    that SELL's own price (the only price we actually have for it).

    This is what "instead of blindly popping up errors" means concretely:
    previously, build_daily_positions clamped the SELL to whatever was
    available (usually 0) and just warned — which silently erased that
    quantity from every downstream calculation (holdings, P&L, attribution,
    the daily ledger), and was the actual cause of daily_positions
    permanently disagreeing with current_positions for any portfolio with
    this pattern. Backfilling the missing opening lot up front means every
    downstream function sees a complete, non-negative position history
    without needing special-case handling anywhere else.

    Pricing the synthetic BUY at the same price as the triggering SELL is a
    deliberate, explainable choice, not a guess dressed up as data: it
    assumes zero gain/loss on the unknown-history portion (we have no idea
    what was actually paid for it), rather than inventing a cost basis.

    Returns (augmented_df, backfill_summary) — the summary is for a clear,
    non-alarming informational message rather than a bare warning.
    """
    empty_summary = pd.DataFrame(
        columns=["portfolio", "ticker_yf", "backfilled_qty", "as_of_date", "price_used", "triggered_by_sell_on"]
    )
    if df.empty:
        return df, empty_summary

    df_sorted = df.sort_values(["portfolio", "ticker_yf", "date"], kind="mergesort").reset_index(drop=True)
    backfill_rows: list[dict] = []
    summary_rows: list[dict] = []

    for (portfolio, ticker), grp in df_sorted.groupby(["portfolio", "ticker_yf"], sort=False):
        running = 0.0
        for _, row in grp.iterrows():
            qty = float(row["quantity"])
            action = str(row["action"]).upper()
            if action == "BUY":
                running += qty
            elif action == "SELL":
                if qty > running + 1e-9:
                    shortfall = qty - running
                    sell_price = float(row["price"])
                    backfill_date = row["date"] - pd.Timedelta(days=1)
                    backfill_rows.append(
                        {
                            "portfolio": portfolio,
                            "date": backfill_date,
                            "ticker": ticker,
                            "ticker_yf": ticker,
                            "action": "BUY",
                            "quantity": shortfall,
                            "price": sell_price,
                            "signed_qty": shortfall,
                            "is_opening": False,
                            "opening_market_value": np.nan,
                        }
                    )
                    summary_rows.append(
                        {
                            "portfolio": portfolio,
                            "ticker_yf": ticker,
                            "backfilled_qty": round(shortfall, 4),
                            "as_of_date": backfill_date.date().isoformat(),
                            "price_used": sell_price,
                            "triggered_by_sell_on": row["date"].date().isoformat(),
                        }
                    )
                    running = 0.0
                else:
                    running -= qty

    if not backfill_rows:
        return df, empty_summary

    backfill_df = pd.DataFrame(backfill_rows)
    augmented = pd.concat([df, backfill_df], ignore_index=True, sort=False)
    augmented = augmented.sort_values(["portfolio", "date", "ticker_yf"], kind="mergesort").reset_index(drop=True)
    return augmented, pd.DataFrame(summary_rows)


def missing_price_tickers(daily_positions: dict[str, pd.DataFrame], prices: pd.DataFrame) -> set[str]:
    """Tickers that appear in someone's holdings but have no downloaded price series."""
    missing: set[str] = set()
    for holdings in daily_positions.values():
        missing.update(set(holdings.columns) - set(prices.columns))
    return missing


def compute_pnl(df: pd.DataFrame, prices: pd.DataFrame) -> pd.DataFrame:
    """
    Compute realized and unrealized P&L per portfolio using FIFO cost basis.

    Returns a DataFrame with columns:
        portfolio, realized_pnl, unrealized_pnl, total_pnl
    """
    results = []

    for portfolio in sorted(df["portfolio"].unique().tolist()):
        txns = df[df["portfolio"] == portfolio].sort_values("date").copy()
        # FIFO queue per ticker: list of (qty, cost_price) lots
        lots: dict[str, list[tuple[float, float]]] = {}
        realized = 0.0

        for _, row in txns.iterrows():
            ticker = str(row["ticker_yf"])
            qty = float(row["quantity"])
            price = float(row["price"])
            action = str(row["action"]).upper()

            if action == "BUY":
                lots.setdefault(ticker, []).append((qty, price))
            elif action == "SELL":
                remaining_sell = qty
                ticker_lots = lots.get(ticker, [])
                new_lots = []
                for lot_qty, lot_price in ticker_lots:
                    if remaining_sell <= 0:
                        new_lots.append((lot_qty, lot_price))
                        continue
                    if lot_qty <= remaining_sell:
                        realized += lot_qty * (price - lot_price)
                        remaining_sell -= lot_qty
                    else:
                        realized += remaining_sell * (price - lot_price)
                        new_lots.append((lot_qty - remaining_sell, lot_price))
                        remaining_sell = 0
                lots[ticker] = new_lots

        # Unrealized: latest price vs avg cost for open lots
        unrealized = 0.0
        last_date = prices.index[-1] if not prices.empty else None
        for ticker, ticker_lots in lots.items():
            if not ticker_lots:
                continue
            open_qty = sum(q for q, _ in ticker_lots)
            if open_qty <= 0:
                continue
            # Get latest market price
            mkt_price = None
            if ticker in prices.columns and last_date is not None:
                s = prices[ticker].dropna()
                if not s.empty:
                    mkt_price = float(s.iloc[-1])
            if mkt_price is None:
                continue
            avg_cost = sum(q * p for q, p in ticker_lots) / open_qty
            unrealized += open_qty * (mkt_price - avg_cost)

        results.append({
            "portfolio": portfolio,
            "realized_pnl": round(realized, 2),
            "unrealized_pnl": round(unrealized, 2),
            "total_pnl": round(realized + unrealized, 2),
        })

    return pd.DataFrame(results).set_index("portfolio")


def drawdown_series(value_series: pd.Series) -> pd.Series:
    """
    Generic peak-to-trough drawdown of any monotonic-ish value series.

    Feed this a TWR growth-of-₹1 curve (`cumulative_growth_from_twr`) for a
    performance drawdown — never raw NAV, which would register a client
    withdrawal as a "drawdown" even when nothing was actually lost.
    """
    value_series = value_series.dropna()
    if value_series.empty:
        return pd.Series(dtype=float)
    rolling_max = value_series.cummax()
    return value_series / rolling_max - 1


def performance_attribution(holdings: pd.DataFrame, prices: pd.DataFrame) -> pd.Series:
    """
    Per-ticker contribution to total portfolio return: for each day, the
    prior day's portfolio weight of a ticker times that ticker's return that
    day, summed over the whole period.

        contribution_i = sum_t( weight_i(t-1) * stock_return_i(t) )

    This is standard daily-rebalanced ("arithmetic") return attribution: the
    sum of all tickers' contributions approximates total portfolio return
    (it won't match exactly — compounding cross-terms are the difference —
    but it should be close; a large gap signals a data problem, e.g. a
    ticker with a bad price series).

    `holdings` must already be non-negative (see build_daily_positions,
    which clips at zero) — a phantom negative holding would flip the sign
    of that ticker's weight and silently corrupt every other ticker's
    contribution too, since weights are normalized by total portfolio value.
    """
    tickers = [t for t in holdings.columns if t in prices.columns]
    if not tickers:
        return pd.Series(dtype=float)

    portfolio_prices = prices[tickers]
    stock_returns = portfolio_prices.pct_change().fillna(0)

    holdings_aligned = holdings[tickers].reindex(stock_returns.index).ffill().fillna(0)
    # Market value of the portfolio's price-covered holdings, day by day —
    # used only to normalize weights below, never displayed as "NAV".
    market_value = (holdings_aligned * portfolio_prices).sum(axis=1)
    weights = holdings_aligned.multiply(portfolio_prices).div(market_value, axis=0).fillna(0)

    contribution = (weights.shift(1).fillna(0) * stock_returns).sum().sort_values(ascending=False)
    return contribution


def attribution_reconciliation(
    holdings: pd.DataFrame, prices: pd.DataFrame, twr_returns: pd.Series
) -> dict[str, float]:
    """
    Sanity check for performance_attribution: compares the sum of all
    per-ticker contributions against the portfolio's actual cumulative TWR
    return over the same window. A small gap is expected (arithmetic vs.
    compounded return); a large one usually means a price or holdings issue.
    """
    attrib = performance_attribution(holdings, prices)
    contribution_total = float(attrib.sum()) if not attrib.empty else np.nan

    r = pd.to_numeric(twr_returns, errors="coerce").dropna()
    actual_total = float((1 + r).prod() - 1) if len(r) >= 1 else np.nan

    gap = (
        float(contribution_total - actual_total)
        if pd.notna(contribution_total) and pd.notna(actual_total)
        else np.nan
    )
    return {
        "contribution_total": contribution_total,
        "actual_cumulative_return": actual_total,
        "gap": gap,
    }


def overlap_matrix(positions_df: pd.DataFrame) -> pd.DataFrame:
    portfolios = sorted(positions_df["portfolio"].unique().tolist())
    om = pd.DataFrame(index=portfolios, columns=portfolios, dtype=float)

    sets: dict[str, set[str]] = {}
    for p in portfolios:
        sets[p] = set(positions_df[positions_df["portfolio"] == p]["ticker_yf"].dropna().astype(str).tolist())

    for p1 in portfolios:
        for p2 in portfolios:
            set1, set2 = sets[p1], sets[p2]
            denom = len(set1.union(set2))
            om.loc[p1, p2] = round((len(set1.intersection(set2)) / denom) * 100, 1) if denom else 0.0

    return om


def risk_contribution(holdings: pd.DataFrame, prices: pd.DataFrame, min_periods: int = 20) -> pd.Series:
    """
    Each ticker's share of total portfolio risk, computed from its ACTUAL
    day-by-day historical weight — not a single fixed snapshot of today's
    weight applied across the whole history.

    A stock that was a large position during a volatile stretch and is a
    small position today should show up as a bigger risk contributor than
    a stock that's large today but was small (or unheld) during the
    volatile period. The previous version couldn't distinguish those,
    since it multiplied a static covariance matrix by a single fixed
    weight vector (today's weights) — that answers "if today's allocation
    had applied the whole time", not "what actually drove risk".

    Method: for each day t, weight_i,t-1 = (qty_i,t-1 * price_i,t-1) /
    portfolio_value_t-1 (yesterday's weight, since a weight is a
    start-of-period allocation — same convention performance_attribution
    uses). Each stock's daily contribution to portfolio return is
    contrib_i,t = weight_i,t-1 * return_i,t, so portfolio_return_t =
    sum_i contrib_i,t. Each stock's risk-contribution share is
    Cov(contrib_i, portfolio_return) / Var(portfolio_return) — a Euler
    decomposition that sums to 1.0 across stocks by construction (renormalized
    defensively in case a ticker gets excluded for insufficient history) and
    correctly weights each period by how much that stock actually mattered
    at the time. A negative share means that stock's moves tended to offset
    the rest of the portfolio (a diversifier reducing risk), which is real
    and kept as-is rather than floored at zero.

    Includes any ticker ever held (not just currently-held), since a
    position that was fully exited can still have been a major historical
    risk driver — that's exactly the kind of thing a fixed-current-weight
    version couldn't show.
    """
    if holdings.empty:
        return pd.Series(dtype=float)

    tickers = [t for t in holdings.columns if t in prices.columns]
    if len(tickers) < 2:
        return pd.Series(dtype=float)

    px = prices[tickers].reindex(holdings.index).ffill()
    qty = holdings[tickers]
    market_value = qty * px
    portfolio_value = market_value.sum(axis=1)

    weights_prev = market_value.shift(1).div(portfolio_value.shift(1).replace(0, np.nan), axis=0)
    stock_returns = px.pct_change()
    contrib = weights_prev * stock_returns  # each stock's daily weighted contribution to portfolio return

    port_return = contrib.sum(axis=1, min_count=1)
    port_var = port_return.var(ddof=1)
    if not port_var or pd.isna(port_var) or port_var <= 0:
        return pd.Series(dtype=float)

    rc_shares: dict[str, float] = {}
    for t in tickers:
        c = contrib[t]
        valid = c.notna() & port_return.notna()
        if valid.sum() < min_periods:
            continue
        cov_i = c[valid].cov(port_return[valid])
        if pd.notna(cov_i):
            rc_shares[t] = cov_i / port_var

    if not rc_shares:
        return pd.Series(dtype=float)

    rc = pd.Series(rc_shares)
    total = rc.sum()
    if total and abs(total) > 1e-9:
        rc = rc / total  # guards against drift from tickers excluded by min_periods
    return rc.sort_values(ascending=False)


def categorize_market_caps_inr(
    market_caps_inr: pd.Series,
    *,
    large_crore: float = 100000,
    mid_crore: float = 30000,
) -> pd.Series:
    """
    Categorize holdings by market-cap bucket.

    ETFs are labelled "ETF" because an ETF's own market cap is not comparable
    to a listed company's equity market cap. REITs retain a separate "REIT"
    bucket only when their actual market cap is unavailable.
    """
    caps = pd.to_numeric(market_caps_inr, errors="coerce")
    large = large_crore * 1e7
    mid = mid_crore * 1e7

    def _cat(ticker: str, value: float) -> str:
        asset_class = ASSET_CLASS_MAP.get(str(ticker).upper())
        if asset_class == "ETF":
            return "ETF"
        if asset_class == "REIT" and pd.isna(value):
            return "REIT"
        if pd.isna(value):
            return "Unknown"
        if value >= large:
            return "Large Cap"
        if value >= mid:
            return "Mid Cap"
        return "Small Cap"

    return pd.Series(
        {ticker: _cat(ticker, value) for ticker, value in caps.items()},
        dtype="object",
    )


def cap_split_weights(
    holdings: pd.DataFrame,
    prices: pd.DataFrame,
    market_caps_inr: pd.Series,
    *,
    large_crore: float = 100000,
    mid_crore: float = 30000,
) -> pd.Series:
    """
    Returns market-value weights by Large/Mid/Small Cap plus separate ETF/REIT
    buckets, so non-company instruments are never dumped into "Unknown".
    """
    tickers = [t for t in holdings.columns if t in prices.columns]
    if not tickers:
        return pd.Series(dtype=float)

    latest_h = holdings.iloc[-1][tickers]
    latest_p = prices.loc[prices.index[-1], tickers]
    mv = (latest_h * latest_p).dropna()
    mv = mv[mv > 0]
    if mv.empty:
        return pd.Series(dtype=float)

    cats = categorize_market_caps_inr(
        market_caps_inr.reindex(mv.index),
        large_crore=large_crore,
        mid_crore=mid_crore,
    )
    split = mv.groupby(cats).sum()
    return (split / split.sum()).sort_values(ascending=False)


# ============================== SECTION: pdf (report generation) ==============================

# Brand palette sampled directly from the VIKA Wealth logo (assets/Vika_Logo.jpg):
# indigo background #290084, gold wordmark #F6BF02, green bar #1D9951, red bar #AA1437.
# Variable names kept as-is (DARK_TEAL, ACCENT_TEAL, ...) since they're referenced
# throughout the PDF-generation code — only the underlying hex values changed.
DARK_TEAL = colors.HexColor("#290084")     # brand indigo — headers, table header bands
ACCENT_TEAL = colors.HexColor("#F6BF02")   # brand gold — divider lines, accents
WHITE = colors.white
LIGHT_GRAY = colors.HexColor("#F2F3F4")
MID_GRAY = colors.HexColor("#717D7E")
BLACK = colors.black
POSITIVE_GREEN = colors.HexColor("#1D9951")  # brand green (logo bar)
NEGATIVE_RED = colors.HexColor("#AA1437")    # brand red (logo bar)


PAGE_W, PAGE_H = A4


def _style(name: str, **kw) -> ParagraphStyle:
    return ParagraphStyle(name, **kw)


BASE = getSampleStyleSheet()


def _ensure_unicode_fonts_registered() -> None:
    """
    Register DejaVu fonts (shipped with matplotlib) so unicode symbols like ₹ render correctly.
    """
    try:
        pdfmetrics.getFont("DejaVuSans")
        return
    except Exception:
        pass

    try:
        from matplotlib import font_manager

        regular = font_manager.findfont("DejaVu Sans")
        bold = font_manager.findfont("DejaVu Sans:style=normal:weight=bold")
        pdfmetrics.registerFont(TTFont("DejaVuSans", regular))
        pdfmetrics.registerFont(TTFont("DejaVuSans-Bold", bold))
    except Exception:
        # If registration fails, fall back to core fonts.
        return


_ensure_unicode_fonts_registered()

TITLE_STYLE = _style(
    "ReportTitle",
    fontName="DejaVuSans-Bold" if "DejaVuSans-Bold" in pdfmetrics.getRegisteredFontNames() else "Helvetica-Bold",
    fontSize=24,
    textColor=DARK_TEAL,
    alignment=TA_CENTER,
    leading=28,
)
SUBTITLE_STYLE = _style(
    "ReportSubtitle",
    fontName="DejaVuSans" if "DejaVuSans" in pdfmetrics.getRegisteredFontNames() else "Helvetica-Oblique",
    fontSize=11,
    textColor=MID_GRAY,
    alignment=TA_CENTER,
    leading=16,
)
SECTION_STYLE = _style(
    "SectionHead",
    fontName="DejaVuSans-Bold" if "DejaVuSans-Bold" in pdfmetrics.getRegisteredFontNames() else "Helvetica-Bold",
    fontSize=12,
    textColor=DARK_TEAL,
    spaceBefore=10,
    spaceAfter=4,
)
BODY_STYLE = _style(
    "Body",
    fontName="DejaVuSans" if "DejaVuSans" in pdfmetrics.getRegisteredFontNames() else "Helvetica",
    fontSize=9,
    textColor=BLACK,
    leading=13,
)


def fmt_pct(v) -> str:
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "N/A"
    return f"{v*100:+.2f}%"


def fmt_inr(val: float) -> str:
    if val is None or (isinstance(val, float) and np.isnan(val)):
        return "N/A"
    return format_inr(val) or "N/A"


def make_chart_image(fig, width_mm: float = 170, height_mm: float = 70) -> Image:
    import traceback as _tb
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    plt.close(fig)
    return Image(buf, width=width_mm * mm, height=height_mm * mm)


def growth_chart_figure(
    returns: pd.Series, name: str, benchmark_returns: Optional[dict[str, pd.Series]] = None
):
    """Builds and returns the raw matplotlib figure — callers decide how
    to embed it (make_chart_image() for the PDF, saved as PNG for pptx)."""
    fig, ax = plt.subplots(figsize=(9, 5.2))
    cumret = (1 + returns).cumprod()
    ax.plot(cumret.index, cumret.values, color="#290084", linewidth=2.0, label=name)
    bench_colors = ["#E67E22", "#8C5A0B"]
    for i, (bname, b_rets) in enumerate((benchmark_returns or {}).items()):
        bench = (1 + b_rets.reindex(returns.index).dropna()).cumprod()
        if bench.empty:
            continue
        ax.plot(
            bench.index, bench.values,
            color=bench_colors[i % len(bench_colors)], linewidth=1.4, linestyle="--", label=bname,
        )
    ax.set_title(f"Performance vs. Benchmarks — {name}", fontsize=11, color="#290084", fontweight="bold")
    ax.set_ylabel("Growth of ₹1", fontsize=9)
    ax.legend(fontsize=8)
    ax.tick_params(labelsize=8)
    fig.tight_layout()
    return fig


def growth_chart(
    returns: pd.Series, name: str, benchmark_returns: Optional[dict[str, pd.Series]] = None
) -> Image:
    fig = growth_chart_figure(returns, name, benchmark_returns=benchmark_returns)
    return make_chart_image(fig, height_mm=100)


def drawdown_chart(returns: pd.Series, name: str) -> Image:
    """Drawdown of the TWR growth-of-₹1 curve — never raw NAV, so a client
    withdrawal never shows up as a fake loss."""
    cumret = (1 + returns).cumprod()
    rolling_max = cumret.cummax()
    dd = cumret / rolling_max - 1
    fig, ax = plt.subplots(figsize=(9, 2.8))
    ax.fill_between(dd.index, dd.values, 0, color="#AA1437", alpha=0.5)
    ax.plot(dd.index, dd.values, color="#AA1437", linewidth=0.8)
    ax.set_title(f"Drawdown — {name}", fontsize=9, color="#290084", fontweight="bold")
    ax.yaxis.set_major_formatter(mticker.PercentFormatter(xmax=1, decimals=0))
    ax.grid(False)
    ax.tick_params(labelsize=7)
    fig.tight_layout()
    return make_chart_image(fig, height_mm=55)


def top_holdings_table(holdings: pd.DataFrame, prices: pd.DataFrame, n: int = 12) -> Optional[Table]:
    tickers = [t for t in holdings.columns if t in prices.columns]
    if not tickers:
        return None

    latest_h = holdings.iloc[-1][tickers]
    latest_p = prices.loc[prices.index[-1], tickers]
    mv = (latest_h * latest_p).sort_values(ascending=False).head(n)
    total = mv.sum()

    rows = [["Ticker", "Market Value (INR)", "Weight"]]
    for ticker, val in mv.items():
        weight_str = f"{val/total*100:.1f}%" if total and pd.notna(total) and abs(total) > 1e-9 else "N/A"
        rows.append([ticker.replace(".NS", "").replace(".BO", ""), fmt_inr(_safe_float(val) or 0.0), weight_str])

    t = Table(rows, colWidths=[70 * mm, 65 * mm, 30 * mm])
    t.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), DARK_TEAL),
                ("TEXTCOLOR", (0, 0), (-1, 0), WHITE),
                ("FONTNAME", (0, 0), (-1, 0), TITLE_STYLE.fontName),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("FONTNAME", (0, 1), (-1, -1), BODY_STYLE.fontName),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [WHITE, LIGHT_GRAY]),
                ("GRID", (0, 0), (-1, -1), 0.3, MID_GRAY),
                ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    return t


def sector_exposure_weights(
    holdings: pd.DataFrame, prices: pd.DataFrame, sector_map: dict[str, str]
) -> pd.Series:
    """Market-value-weighted sector exposure (fraction of total, sums to
    1.0) for the latest holdings — shared by both the live UI chart and the
    PDF report, so a ticker's classification is never inconsistent between
    the two."""
    weights, _ = sector_exposure_weights_and_counts(holdings, prices, sector_map)
    return weights


def sector_exposure_weights_and_counts(
    holdings: pd.DataFrame, prices: pd.DataFrame, sector_map: dict[str, str]
) -> tuple[pd.Series, pd.Series]:
    """Same as sector_exposure_weights, plus the number of distinct
    currently-held tickers in each sector — used for pie labels like
    'Financials (34.2%, 3 stocks)'."""
    tickers = [t for t in holdings.columns if t in prices.columns]
    if not tickers:
        return pd.Series(dtype=float), pd.Series(dtype=int)
    latest_h = holdings.iloc[-1][tickers]
    latest_p = prices.loc[prices.index[-1], tickers]
    mv = latest_h * latest_p
    mv = mv[mv > 0]
    if mv.empty:
        return pd.Series(dtype=float), pd.Series(dtype=int)
    sector_of = pd.Series(mv.index, index=mv.index).map(sector_map).fillna("Unknown")
    exp = mv.groupby(sector_of).sum().sort_values(ascending=False)
    counts = sector_of.value_counts().reindex(exp.index)
    total = exp.sum()
    weights = exp / total if total else exp
    return weights, counts


def _make_autopct_with_counts(counts: pd.Series):
    """Builds an autopct callback for ax.pie()/Series.plot.pie() that shows
    the percentage plus a stock count in parentheses, e.g. '34.2% (3)'.
    matplotlib invokes autopct once per wedge in the same order as the
    input data, so a simple counter closure over `counts` (already
    reindexed to match that same order) lines each call up with the right
    sector correctly."""
    values = counts.tolist()
    state = {"i": 0}

    def _autopct(pct):
        idx = state["i"]
        n = int(values[idx]) if idx < len(values) else 0
        state["i"] += 1
        return f"{pct:.1f}% ({n})"

    return _autopct


def sector_pie_figure_clean(exp: pd.Series):
    """
    Clean pie chart, wedges only — no on-slice labels, no autopct text.
    A pie with many small categories (the sector breakdown routinely has
    15-20+) gets unreadably cluttered with labels crammed around the rim;
    the actual numbers belong in the legend table next to it instead,
    where there's room to lay them out properly.

    Returns (fig, colors_used) — colors_used is the exact color assigned
    to each wedge in `exp`'s order, so a companion legend/table can use
    swatches that actually match.
    """
    if exp.empty:
        return None, []
    wedge_colors = list(plt.cm.tab20.colors) + list(plt.cm.tab20b.colors)
    colors_used = [wedge_colors[i % len(wedge_colors)] for i in range(len(exp))]
    fig, ax = plt.subplots(figsize=(6, 6))
    ax.pie(exp.values, colors=colors_used, startangle=90, wedgeprops={"edgecolor": "white", "linewidth": 1})
    ax.set_title("Sector Exposure", fontsize=11, color="#290084", fontweight="bold")
    fig.tight_layout()
    colors_hex = [matplotlib.colors.to_hex(c) for c in colors_used]
    return fig, colors_hex


def sector_pie_image(holdings: pd.DataFrame, prices: pd.DataFrame, sector_map: dict[str, str]):
    exp, counts = sector_exposure_weights_and_counts(holdings, prices, sector_map)
    if exp.empty:
        return None

    fig, colors_hex = sector_pie_figure_clean(exp)
    pie_img = make_chart_image(fig, width_mm=90, height_mm=90)

    rows = [["", "Sector", "%", "Stocks"]]
    for sector, weight in exp.items():
        rows.append(["", sector, f"{weight*100:.1f}%", str(int(counts.get(sector, 0)))])

    legend_tbl = Table(rows, colWidths=[6 * mm, 45 * mm, 16 * mm, 16 * mm], rowHeights=[7 * mm] + [6 * mm] * len(exp))
    style = [
        ("BACKGROUND", (0, 0), (-1, 0), DARK_TEAL),
        ("TEXTCOLOR", (1, 0), (-1, 0), WHITE),
        ("FONTNAME", (1, 0), (-1, 0), TITLE_STYLE.fontName),
        ("FONTNAME", (0, 1), (-1, -1), BODY_STYLE.fontName),
        ("FONTSIZE", (0, 0), (-1, -1), 7.5),
        ("ROWBACKGROUNDS", (1, 1), (-1, -1), [WHITE, LIGHT_GRAY]),
        ("GRID", (1, 0), (-1, -1), 0.3, MID_GRAY),
        ("ALIGN", (2, 1), (-1, -1), "RIGHT"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 2),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
    ]
    for i, hexcolor in enumerate(colors_hex, start=1):
        style.append(("BACKGROUND", (0, i), (0, i), colors.HexColor(hexcolor)))
    legend_tbl.setStyle(TableStyle(style))

    outer = Table([[pie_img, legend_tbl]], colWidths=[95 * mm, 85 * mm])
    outer.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "MIDDLE")]))
    return outer


def stock_risk_return(attribution: pd.Series, risk_contrib: pd.Series) -> pd.DataFrame:
    """
    Combines each stock's contribution to TOTAL portfolio return
    (performance_attribution) and contribution to TOTAL portfolio risk
    (risk_contribution) into one table — this is deliberately NOT each
    stock's own standalone annualized return/volatility. Both inputs
    already sum to 100% across all holdings, so this shows exactly how
    much of the whole return (or risk) pie each stock is responsible for
    — a small, low-volatility position that's heavily overweighted can
    still be a large risk contributor, and that's the point of plotting
    contribution rather than each stock's isolated stats.

    Returns columns [ticker, return_contribution, risk_contribution], one
    row per ticker that appears in either series.
    """
    idx = sorted(set(attribution.index) | set(risk_contrib.index))
    if not idx:
        return pd.DataFrame(columns=["ticker", "return_contribution", "risk_contribution"])
    df = pd.DataFrame(index=idx)
    df["return_contribution"] = attribution.reindex(idx)
    df["risk_contribution"] = risk_contrib.reindex(idx)
    df = df.dropna(how="all").fillna(0.0)

    # performance_attribution's raw output is in percentage-POINTS of the
    # actual period return (sums to the period return, e.g. 12.4%, not
    # 100%) — that's the right convention for the separate "Contribution
    # to Return" table elsewhere, which people read as "this stock added
    # 2.8 points to my return". For THIS chart specifically, normalize it
    # into a share of the total return (sums to 100%, same convention as
    # risk_contribution) so both axes genuinely answer "what % of the
    # whole pie". A stock with a contribution of the opposite sign to the
    # total return correctly comes out negative here (it dragged against
    # the overall result), same interpretation as a negative risk share.
    total_return = df["return_contribution"].sum()
    if total_return and abs(total_return) > 1e-9:
        df["return_contribution"] = df["return_contribution"] / total_return

    return df.reset_index().rename(columns={"index": "ticker"})


def risk_return_quadrants(rr_df: pd.DataFrame) -> pd.DataFrame:
    """
    Labels each stock's quadrant relative to the MEDIAN risk contribution
    and MEDIAN return contribution across the plotted set — "high" and
    "low" are relative to this specific portfolio's own stocks. Since
    contribution values can be negative (a stock that lost money, or one
    whose moves offset the rest of the book and reduced risk), the median
    itself can sit below zero — that's expected and still a meaningful
    split for this portfolio's own holdings.
    """
    if rr_df.empty:
        return rr_df.assign(quadrant=pd.Series(dtype=str))
    risk_mid = rr_df["risk_contribution"].median()
    ret_mid = rr_df["return_contribution"].median()
    def _label(row):
        risk = "High Risk Contribution" if row["risk_contribution"] >= risk_mid else "Low Risk Contribution"
        ret = "High Return Contribution" if row["return_contribution"] >= ret_mid else "Low Return Contribution"
        return f"{risk} / {ret}"
    out = rr_df.copy()
    out["quadrant"] = out.apply(_label, axis=1)
    out.attrs["risk_mid"] = risk_mid
    out.attrs["ret_mid"] = ret_mid
    return out


def risk_return_scatter_figure(rr_df: pd.DataFrame):
    """Matplotlib figure: % contribution to total portfolio risk (x) vs %
    contribution to total portfolio return (y) per stock — both axes sum
    to 100% across all holdings — quartered at the median of each. Returns
    None if there isn't enough data. Caller decides how to render/embed
    the figure (st.pyplot for the live UI, make_chart_image for the PDF).

    Decluttered vs. a plain scatter: each quadrant gets a light background
    tint and its points share that quadrant's color (with a legend), so
    the picture reads at a glance without having to parse every label
    individually — the labels are still there for the exact ticker, but
    color + position carries the story on their own. Labels get a
    semi-transparent white backing so they stay legible over overlapping
    points, and a bigger figure gives everything more room to breathe.
    """
    if rr_df is None or rr_df.empty:
        return None
    labeled = risk_return_quadrants(rr_df)
    x = labeled["risk_contribution"] * 100
    y = labeled["return_contribution"] * 100
    risk_mid = labeled.attrs.get("risk_mid", x.median() / 100) * 100
    ret_mid = labeled.attrs.get("ret_mid", y.median() / 100) * 100

    quadrant_colors = {
        "High Risk Contribution / High Return Contribution": "#F6BF02",   # gold
        "High Risk Contribution / Low Return Contribution": "#AA1437",    # red
        "Low Risk Contribution / High Return Contribution": "#1D9951",    # green
        "Low Risk Contribution / Low Return Contribution": "#7B8794",     # neutral gray
    }
    quadrant_short = {
        "High Risk Contribution / High Return Contribution": "High risk / High return",
        "High Risk Contribution / Low Return Contribution": "High risk / Low return",
        "Low Risk Contribution / High Return Contribution": "Low risk / High return",
        "Low Risk Contribution / Low Return Contribution": "Low risk / Low return",
    }

    x_pad = max((x.max() - x.min()) * 0.25, 5)
    y_pad = max((y.max() - y.min()) * 0.25, 5)
    x_lo, x_hi = x.min() - x_pad, x.max() + x_pad
    y_lo, y_hi = y.min() - y_pad, y.max() + y_pad

    fig, ax = plt.subplots(figsize=(8.5, 7))

    # Quadrant background tints, drawn first so points/labels sit on top.
    ax.axvspan(risk_mid, x_hi, ymin=(ret_mid - y_lo) / (y_hi - y_lo), ymax=1, color="#F6BF02", alpha=0.06, zorder=0)
    ax.axvspan(risk_mid, x_hi, ymin=0, ymax=(ret_mid - y_lo) / (y_hi - y_lo), color="#AA1437", alpha=0.06, zorder=0)
    ax.axvspan(x_lo, risk_mid, ymin=(ret_mid - y_lo) / (y_hi - y_lo), ymax=1, color="#1D9951", alpha=0.06, zorder=0)
    ax.axvspan(x_lo, risk_mid, ymin=0, ymax=(ret_mid - y_lo) / (y_hi - y_lo), color="#7B8794", alpha=0.06, zorder=0)

    for _, row in labeled.iterrows():
        color = quadrant_colors.get(row["quadrant"], "#290084")
        ax.scatter(
            row["risk_contribution"] * 100, row["return_contribution"] * 100,
            s=90, color=color, zorder=3, edgecolors="white", linewidths=0.8,
        )
        label = str(row["ticker"]).replace(".NS", "").replace(".BO", "")
        ax.annotate(
            label, (row["risk_contribution"] * 100, row["return_contribution"] * 100),
            fontsize=8, xytext=(6, 5), textcoords="offset points", color="#290084",
            bbox=dict(boxstyle="round,pad=0.15", facecolor="white", edgecolor="none", alpha=0.65),
        )

    ax.axvline(risk_mid, color="#999999", linestyle="--", linewidth=0.8, zorder=1)
    ax.axhline(ret_mid, color="#999999", linestyle="--", linewidth=0.8, zorder=1)
    ax.axvline(0, color="#cccccc", linewidth=0.6, zorder=0)
    ax.axhline(0, color="#cccccc", linewidth=0.6, zorder=0)
    ax.set_xlim(x_lo, x_hi)
    ax.set_ylim(y_lo, y_hi)
    ax.set_xlabel("% Contribution to Total Portfolio Risk", fontsize=10)
    ax.set_ylabel("% Share of Total Portfolio Return", fontsize=10)
    ax.set_title("Risk vs. Return Contribution by Holding", fontsize=11, color="#290084", fontweight="bold")
    ax.tick_params(labelsize=9)
    ax.grid(False)

    legend_handles = [
        plt.Line2D([0], [0], marker="o", color="w", markerfacecolor=c, markeredgecolor="white", markersize=9, label=quadrant_short[q])
        for q, c in quadrant_colors.items()
    ]
    ax.legend(handles=legend_handles, loc="upper center", bbox_to_anchor=(0.5, -0.12), ncol=2, fontsize=8, frameon=False)

    fig.tight_layout()
    return fig


def _contrast_text_color(hex_color: str) -> str:
    """White text on a dark wedge/cell, black text on a light one — plain
    luminance threshold, avoids illegible same-hue-on-same-hue labels
    (e.g. dark navy percentage text on a solid indigo pie wedge)."""
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return "#FFFFFF" if luminance < 140 else "#1B1B1B"


def cap_split_pie_image(cap_split: pd.Series) -> Optional[Image]:
    """Market cap allocation pie (Large/Mid/Small Cap) — cap_split is a
    Series of bucket -> weight (0-1), e.g. from portfolio_pipeline's
    cap_split_weights()."""
    s = pd.to_numeric(cap_split, errors="coerce").dropna()
    s = s[s > 0]
    if s.empty:
        return None

    color_map = {"Large Cap": "#290084", "Mid Cap": "#F6BF02", "Small Cap": "#1D9951", "ETF": "#717D7E", "REIT": "#AA1437", "Unknown": "#B0B0B0"}
    colors_ordered = [color_map.get(k, "#B0B0B0") for k in s.index]

    fig, ax = plt.subplots(figsize=(5, 5))
    _, _, autotexts = ax.pie(
        s.values, labels=s.index, autopct="%1.1f%%", startangle=90,
        colors=colors_ordered, textprops={"fontsize": 8},
    )
    for txt, color in zip(autotexts, colors_ordered):
        txt.set_color(_contrast_text_color(color))
    ax.set_title("Market Cap Allocation", fontsize=9, color="#290084", fontweight="bold")
    fig.tight_layout()
    return make_chart_image(fig, width_mm=90, height_mm=85)


def bar_series_image(series: pd.Series, title: str, *, top_n: int = 12, as_percent: bool = True) -> Optional[Image]:
    s = pd.to_numeric(series, errors="coerce").dropna()
    if s.empty:
        return None
    s = s.head(top_n)
    if as_percent:
        s = s * 100
    fig, ax = plt.subplots(figsize=(9, 3.2))
    s.sort_values(ascending=True).plot.barh(ax=ax, color="#F6BF02")
    if as_percent:
        ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"{v:.1f}%"))
    ax.set_title(title, fontsize=9, color="#290084", fontweight="bold")
    ax.grid(False)
    ax.tick_params(labelsize=7)
    fig.tight_layout()
    return make_chart_image(fig, height_mm=70)


def trailing_returns_pdf_table(trailing_table: pd.DataFrame) -> Optional[Table]:
    """PDF rendering of trailing_returns_table()'s output — same blank-if-
    unavailable convention as the live UI, not 'N/A' or a dash. Rows are
    Portfolio / benchmark / Alpha, columns are the trailing periods."""
    if trailing_table is None or trailing_table.empty:
        return None

    header = [""] + list(trailing_table.columns)
    rows = [header]
    for label, row in trailing_table.iterrows():
        rows.append(
            [label] + ["" if pd.isna(v) else f"{v*100:.2f}%" for v in row]
        )

    n_cols = len(header)
    col_width = 155 * mm / n_cols
    t = Table(rows, colWidths=[col_width] * n_cols)
    style = [
        ("BACKGROUND", (0, 0), (-1, 0), DARK_TEAL),
        ("TEXTCOLOR", (0, 0), (-1, 0), WHITE),
        ("FONTNAME", (0, 0), (-1, 0), TITLE_STYLE.fontName),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("FONTNAME", (0, 1), (-1, -1), BODY_STYLE.fontName),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [WHITE, LIGHT_GRAY]),
        ("GRID", (0, 0), (-1, -1), 0.3, MID_GRAY),
        ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]
    t.setStyle(TableStyle(style))
    return t


def contribution_table(attribution: pd.Series) -> Optional[Table]:
    """
    Full per-ticker contribution-to-return table, since portfolio inception,
    sorted from the biggest positive contributor to the biggest negative
    contributor. Unlike the old top-12 bar chart, this includes every
    ticker so negative contributors are never cut off. Platypus splits
    long tables across pages automatically, so no truncation is needed.

    Shows both the absolute contribution (percentage-points added to
    total return) and its normalized share of the total (summing to
    100% across all tickers) side by side.
    """
    s = pd.to_numeric(attribution, errors="coerce").dropna().sort_values(ascending=False)
    if s.empty:
        return None
    total = s.sum()

    rows = [["Ticker", "Contribution to Return", "% of Total Return"]]
    for ticker, val in s.items():
        pct_of_total = f"{val/total*100:+.2f}%" if total and abs(total) > 1e-9 else "N/A"
        rows.append([str(ticker).replace(".NS", "").replace(".BO", ""), f"{val*100:+.2f}%", pct_of_total])

    t = Table(rows, colWidths=[65 * mm, 55 * mm, 55 * mm], repeatRows=1)
    style = [
        ("BACKGROUND", (0, 0), (-1, 0), DARK_TEAL),
        ("TEXTCOLOR", (0, 0), (-1, 0), WHITE),
        ("FONTNAME", (0, 0), (-1, 0), TITLE_STYLE.fontName),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("FONTNAME", (0, 1), (-1, -1), BODY_STYLE.fontName),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [WHITE, LIGHT_GRAY]),
        ("GRID", (0, 0), (-1, -1), 0.3, MID_GRAY),
        ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]
    for i, (_, val) in enumerate(s.items(), start=1):
        style.append(("TEXTCOLOR", (1, i), (2, i), POSITIVE_GREEN if val >= 0 else NEGATIVE_RED))
    t.setStyle(TableStyle(style))
    return t


def risk_contribution_table(risk_contrib: pd.Series, portfolio_vol: float) -> Optional[Table]:
    """
    Per-ticker risk-contribution table — the absolute contribution to
    total portfolio annualized volatility (percentage-points, sums back
    to the total) alongside the normalized share of total risk (sums to
    100%). Replaces the old bar-chart-only rendering so the actual
    numbers are readable, not just relative bar lengths.
    """
    s = pd.to_numeric(risk_contrib, errors="coerce").dropna().sort_values(ascending=False)
    if s.empty:
        return None

    rows = [["Ticker", "Contribution to Volatility", "% of Total Risk"]]
    for ticker, share in s.items():
        vol_contrib = f"{share*portfolio_vol*100:+.2f}%" if pd.notna(portfolio_vol) else "N/A"
        rows.append([str(ticker).replace(".NS", "").replace(".BO", ""), vol_contrib, f"{share*100:+.2f}%"])

    t = Table(rows, colWidths=[65 * mm, 55 * mm, 55 * mm], repeatRows=1)
    style = [
        ("BACKGROUND", (0, 0), (-1, 0), DARK_TEAL),
        ("TEXTCOLOR", (0, 0), (-1, 0), WHITE),
        ("FONTNAME", (0, 0), (-1, 0), TITLE_STYLE.fontName),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("FONTNAME", (0, 1), (-1, -1), BODY_STYLE.fontName),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [WHITE, LIGHT_GRAY]),
        ("GRID", (0, 0), (-1, -1), 0.3, MID_GRAY),
        ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]
    for i, (_, share) in enumerate(s.items(), start=1):
        style.append(("TEXTCOLOR", (1, i), (2, i), POSITIVE_GREEN if share >= 0 else NEGATIVE_RED))
    t.setStyle(TableStyle(style))
    return t


def pie_series_image(series: pd.Series, title: str) -> Optional[Image]:
    s = pd.to_numeric(series, errors="coerce").dropna()
    if s.empty:
        return None
    fig, ax = plt.subplots(figsize=(5.5, 4.5))
    colorset = plt.cm.Set3.colors
    s.plot.pie(
        ax=ax,
        autopct="%1.0f%%",
        startangle=90,
        colors=colorset[: len(s)],
        textprops={"fontsize": 8},
    )
    ax.set_title(title, fontsize=9, color="#290084", fontweight="bold")
    ax.set_ylabel("")
    fig.tight_layout()
    return make_chart_image(fig, width_mm=95, height_mm=80)


def correlation_summary(stock_corr: Optional[pd.DataFrame], portfolio_returns: Optional[pd.Series], benchmark_returns: Optional[pd.Series]) -> dict:
    """
    Two single-number summaries to sit under the stock-correlation heatmap:
    - avg_pairwise: average pairwise correlation among the portfolio's own
      holdings (the off-diagonal of stock_corr) — a quick read on internal
      diversification; closer to 0 means the stocks move more independently
      of each other, closer to 1 means they tend to move together.
    - vs_benchmark: correlation between the PORTFOLIO's own daily returns
      (not any single stock) and the benchmark's daily returns.
    Either can come back NaN if there isn't enough data.
    """
    avg_pairwise = np.nan
    if stock_corr is not None and not stock_corr.empty and stock_corr.shape[0] >= 2:
        arr = stock_corr.to_numpy(dtype=float)
        n = arr.shape[0]
        mask = ~np.eye(n, dtype=bool)
        off_diag = arr[mask]
        off_diag = off_diag[~np.isnan(off_diag)]
        if off_diag.size:
            avg_pairwise = float(off_diag.mean())

    vs_benchmark = np.nan
    if portfolio_returns is not None and benchmark_returns is not None:
        p = pd.to_numeric(portfolio_returns, errors="coerce").dropna()
        b = pd.to_numeric(benchmark_returns, errors="coerce").dropna()
        aligned = pd.concat([p.rename("p"), b.rename("b")], axis=1).dropna()
        # Guard against near-zero variance on either side — pandas/numpy's
        # corrcoef divides by each series' std dev, which raises a
        # RuntimeWarning (and returns a meaningless value) on a very short
        # or nearly-constant series rather than a real correlation.
        if len(aligned) >= 3 and aligned["p"].std(ddof=1) > 1e-12 and aligned["b"].std(ddof=1) > 1e-12:
            c = aligned["p"].corr(aligned["b"])
            vs_benchmark = float(c) if pd.notna(c) else np.nan

    return {"avg_pairwise": avg_pairwise, "vs_benchmark": vs_benchmark}


def corr_heatmap_figure(corr: pd.DataFrame, title: str = "Correlation heatmap (stocks)"):
    """Builds and returns the raw matplotlib figure, or None if there's
    not enough data — callers decide how to embed it.

    Only the lower triangle (plus diagonal) is drawn — a correlation
    matrix is symmetric, so the upper triangle is a mirror image of the
    lower one and shows nothing new, just doubles the ink."""
    if corr is None or corr.empty:
        return None
    df = corr.copy()
    df = df.apply(pd.to_numeric, errors="coerce")
    df = df.dropna(axis=0, how="all").dropna(axis=1, how="all")
    if df.shape[0] < 2 or df.shape[1] < 2:
        return None

    cmap = get_corr_cmap()
    cmap = cmap.copy() if hasattr(cmap, "copy") else cmap
    cmap.set_bad(color="white")

    values = df.values.astype(float).copy()
    if values.shape[0] == values.shape[1]:
        upper = np.triu_indices(values.shape[0], k=1)
        values[upper] = np.nan

    fig, ax = plt.subplots(figsize=(8.5, 6))
    im = ax.imshow(values, aspect="auto", cmap=cmap, vmin=-1, vmax=1)
    ax.set_xticks(np.arange(len(df.columns)))
    ax.set_yticks(np.arange(len(df.index)))
    ax.set_xticklabels([str(c).replace(".NS", "") for c in df.columns], rotation=45, ha="right", fontsize=6)
    ax.set_yticklabels([str(i).replace(".NS", "") for i in df.index], fontsize=6)
    for i in range(len(df.index)):
        for j in range(len(df.columns)):
            v = values[i, j]
            if pd.isna(v):
                continue
            r, g, b, _ = cmap((float(v) + 1) / 2)
            luminance = 0.299 * r + 0.587 * g + 0.114 * b
            text_color = "white" if luminance < 0.55 else "black"
            ax.text(j, i, f"{v:.2f}", ha="center", va="center", fontsize=5, color=text_color)
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    ax.set_title(title, fontsize=9, color="#290084", fontweight="bold")
    fig.tight_layout()
    return fig


def corr_heatmap_image(corr: pd.DataFrame, title: str) -> Optional[Image]:
    fig = corr_heatmap_figure(corr, title)
    if fig is None:
        return None
    return make_chart_image(fig, width_mm=170, height_mm=120)


def returns_histogram_figure(returns: pd.Series, title: str = "Histogram of daily returns"):
    """Builds and returns the raw matplotlib figure, or None if there's
    not enough data — callers decide how to embed it."""
    r = pd.to_numeric(returns, errors="coerce").dropna()
    if len(r) < 10:
        return None
    cutoff = float(r.quantile(0.01))
    tail = r[r <= cutoff]
    cvar99_val = float(-tail.mean()) if not tail.empty else float(-cutoff)
    fig, ax = plt.subplots(figsize=(8.5, 5.2))
    counts, bin_edges, _ = ax.hist(
        r.values, bins=40, color="#F6BF02", alpha=0.85, edgecolor="white", linewidth=0.3, label="Daily returns"
    )
    # Fitted normal ("bell") curve, scaled to match histogram COUNTS (not a
    # density curve) so it overlays visually on the same axis as the bars.
    mu, sigma = float(r.mean()), float(r.std(ddof=1))
    if sigma > 0 and len(bin_edges) > 1:
        bin_width = bin_edges[1] - bin_edges[0]
        x_curve = np.linspace(r.min(), r.max(), 200)
        pdf = (1.0 / (sigma * np.sqrt(2 * np.pi))) * np.exp(-0.5 * ((x_curve - mu) / sigma) ** 2)
        y_curve = pdf * len(r) * bin_width
        ax.plot(x_curve, y_curve, color="#290084", linewidth=2, label="Fitted normal distribution")
    ax.axvline(cutoff, color="#AA1437", linewidth=1.5, linestyle="--",
               label=f"99% CVaR: {cvar99_val*100:.2f}% (tail avg)")
    ax.set_title(title, fontsize=10, color="#290084", fontweight="bold")
    ax.set_xlabel("Daily return")
    ax.set_ylabel("Frequency")
    ax.xaxis.set_major_formatter(mticker.PercentFormatter(xmax=1, decimals=1))
    ax.legend(fontsize=8)
    ax.grid(False)
    fig.tight_layout()
    return fig


def returns_histogram_image(returns: pd.Series, title: str) -> Optional[Image]:
    fig = returns_histogram_figure(returns, title)
    if fig is None:
        return None
    return make_chart_image(fig, width_mm=170, height_mm=100)


def returns_stats_table(returns: pd.Series) -> Optional[Table]:
    """
    Small descriptive-statistics table meant to sit directly under the
    return histogram: mean, standard deviation, skewness, (excess)
    kurtosis, min, max. Skewness/kurtosis use pandas' own definitions —
    `.kurt()` is EXCESS kurtosis (0 for a normal distribution), matching
    the fitted bell curve drawn on the histogram, so a reader can compare
    "kurtosis = +3.2" against "does the histogram look fat-tailed" directly.
    """
    r = pd.to_numeric(returns, errors="coerce").dropna()
    if len(r) < 10:
        return None
    rows = [
        ["Statistic", "Value"],
        ["Mean (daily)", f"{float(r.mean())*100:+.3f}%"],
        ["Std. Dev. (daily)", f"{float(r.std(ddof=1))*100:.3f}%"],
        ["Skewness", f"{float(r.skew()):+.2f}"],
        ["Kurtosis (excess)", f"{float(r.kurt()):+.2f}"],
        ["Min (daily)", f"{float(r.min())*100:+.2f}%"],
        ["Max (daily)", f"{float(r.max())*100:+.2f}%"],
    ]
    t = Table(rows, colWidths=[70 * mm, 50 * mm])
    t.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), DARK_TEAL),
                ("TEXTCOLOR", (0, 0), (-1, 0), WHITE),
                ("FONTNAME", (0, 0), (-1, 0), TITLE_STYLE.fontName),
                ("FONTNAME", (0, 1), (-1, -1), BODY_STYLE.fontName),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [WHITE, LIGHT_GRAY]),
                ("GRID", (0, 0), (-1, -1), 0.3, MID_GRAY),
                ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
                ("TOPPADDING", (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    return t


def normalize_cap_split_for_report(cap_split: Optional[pd.Series]) -> Optional[pd.Series]:
    """
    Streamlit app may produce 3-5 buckets; for the PDF keep it clean:
    - Drop Unknown if present
    """
    if cap_split is None:
        return None
    s = pd.to_numeric(cap_split, errors="coerce").dropna()
    if s.empty:
        return None
    s = s.copy()
    if "Unknown" in s.index:
        s = s.drop(index=["Unknown"])
    total = float(s.sum())
    if total > 0:
        s = s / total
    return s.sort_values(ascending=False)


def summary_text(
    *,
    portfolio_name: str,
    stats_row: Optional[pd.Series],
    returns: pd.Series,
) -> str:
    r = pd.to_numeric(returns, errors="coerce").dropna()
    start = r.index.min().date().isoformat() if not r.empty else "N/A"
    end = r.index.max().date().isoformat() if not r.empty else "N/A"
    total_ret = float((1 + r).prod() - 1) if len(r) >= 1 else np.nan

    vol = stats_row.get("Period Volatility (TWR)") if stats_row is not None else np.nan
    sharpe = stats_row.get("Sharpe") if stats_row is not None else np.nan
    sortino = stats_row.get("Sortino") if stats_row is not None else np.nan
    var99 = stats_row.get("CVaR 99% (Daily)") if stats_row is not None else np.nan
    mdd = stats_row.get("Max Drawdown") if stats_row is not None else np.nan

    def _pct(v):
        try:
            f = float(v)
            import math
            if math.isnan(f) or math.isinf(f): return "N/A"
            return f"{f*100:.1f}%"
        except (TypeError, ValueError):
            return "N/A"

    def _flt(v):
        try:
            f = float(v)
            import math
            if math.isnan(f) or math.isinf(f): return "N/A"
            return f"{f:.2f}"
        except (TypeError, ValueError):
            return "N/A"

    beta_txt = ""
    if stats_row is not None:
        bench_names = sorted({
            k[len("Beta ("):-1] for k in stats_row.index if k.startswith("Beta (") and k.endswith(")")
        })
        parts_beta = [f"Beta ({name}): {_flt(stats_row.get(f'Beta ({name})'))}" for name in bench_names]
        beta_txt = "; ".join(parts_beta)

    parts = [
        f"This report summarises {portfolio_name}'s equity portfolio performance from {start} to {end}.",
        f"Over the period, the portfolio returned {_pct(total_ret)}.",
        f"Period volatility (TWR): {_pct(vol)}; Sharpe: {_flt(sharpe)}; Sortino: {_flt(sortino)}"
        + (f"; {beta_txt}" if beta_txt else "")
        + f"; 99% daily CVaR: {_pct(var99)}; max drawdown: {_pct(mdd)}.",
    ]
    return " ".join(parts)



def _safe_float(val):
    """Return float or None if val is missing/invalid."""
    if val is None:
        return None
    try:
        v = float(val)
    except (TypeError, ValueError):
        return None
    import math
    return None if (math.isnan(v) or math.isinf(v)) else v


def _safe_fmt_pct(val) -> str:
    """fmt_pct with None/NaN guard — avoids float(None) crash."""
    v = _safe_float(val)
    return "N/A" if v is None else fmt_pct(v)


def _safe_fmt_float(val, decimals: int = 2) -> str:
    """Format a float to N decimal places safely."""
    v = _safe_float(val)
    return "N/A" if v is None else f"{v:.{decimals}f}"


def generate_portfolio_pdf_bytes(
    *,
    portfolio_name: str,
    report_date: date | datetime,
    returns: pd.Series,
    holdings: pd.DataFrame,
    prices: pd.DataFrame,
    stats_row: Optional[pd.Series] = None,
    benchmark_returns: Optional[dict[str, pd.Series]] = None,
    sector_map: Optional[dict[str, str]] = None,
    risk_contrib: Optional[pd.Series] = None,
    attribution: Optional[pd.Series] = None,
    stock_corr: Optional[pd.DataFrame] = None,
    returns_hist: Optional[pd.Series] = None,
    cash_balance: Optional[float] = None,
    portfolio_value: Optional[float] = None,
    cap_split: Optional[pd.Series] = None,
    trailing_returns: Optional[pd.DataFrame] = None,
) -> bytes:
    import traceback as _tb
    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=14 * mm,
        rightMargin=14 * mm,
        topMargin=14 * mm,
        bottomMargin=14 * mm,
        title=f"Portfolio Report - {portfolio_name}",
    )

    sector_map = sector_map or {}

    story = []
    # Support both .png and .jpg logo filenames
    _assets = Path(__file__).resolve().parent / "assets"
    logo_file = next(
        (p for p in [
            _assets / "vika_logo.png",
            _assets / "vika_logo.jpg",
            _assets / "Vika_Logo.png",
            _assets / "Vika_Logo.jpg",
        ] if p.exists()),
        _assets / "vika_logo.png",  # fallback path (may not exist — caught below)
    )

    def _page_header(canvas_obj, doc_obj):
        if logo_file.exists():
            try:
                w = 42 * mm
                h = 17 * mm
                # TOP-RIGHT corner
                x = PAGE_W - doc_obj.rightMargin - w
                y = PAGE_H - doc_obj.topMargin - h
                canvas_obj.drawImage(
                    str(logo_file),
                    x,
                    y,
                    width=w,
                    height=h,
                    preserveAspectRatio=True,
                    mask="auto",
                )
            except Exception:
                pass

    story.append(Spacer(1, 12 * mm))
    story.append(Paragraph("Portfolio Report", TITLE_STYLE))
    story.append(Spacer(1, 2 * mm))
    story.append(Paragraph(str(portfolio_name), SUBTITLE_STYLE))
    story.append(Spacer(1, 2 * mm))
    story.append(Paragraph(f"Report date: {pd.to_datetime(report_date).date().isoformat()}", SUBTITLE_STYLE))
    story.append(Spacer(1, 8 * mm))
    story.append(HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.7))
    story.append(Spacer(1, 8 * mm))

    story.append(KeepTogether([
        Paragraph("Executive Summary", SECTION_STYLE),
        HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.5),
        Spacer(1, 4 * mm),
        Paragraph(
            summary_text(
                portfolio_name=str(portfolio_name),
                stats_row=stats_row,
                returns=returns,
            ),
            BODY_STYLE,
        ),
    ]))
    story.append(Spacer(1, 6 * mm))

    if stats_row is not None:
        rows = []
        if portfolio_value is not None:
            rows.append(["Current Portfolio Value", fmt_inr(portfolio_value)])
        if cash_balance is not None:
            rows.append(["Cash Balance", fmt_inr(cash_balance)])
        rows += [
            ["Period Return (TWR)",     _safe_fmt_pct(stats_row.get("Period Return (TWR)"))],
            ["Period Volatility (TWR)", _safe_fmt_pct(stats_row.get("Period Volatility (TWR)"))],
            ["Sharpe",                _safe_fmt_float(stats_row.get("Sharpe"))],
            ["Sortino",               _safe_fmt_float(stats_row.get("Sortino"))],
        ]
        bench_names = sorted({
            k[len("Beta ("):-1] for k in stats_row.index if k.startswith("Beta (") and k.endswith(")")
        })
        for bname in bench_names:
            rows.append([f"Beta ({bname})", _safe_fmt_float(stats_row.get(f"Beta ({bname})"))])
            rows.append([f"Jensen Alpha ({bname}) (TWR trailing period)", _safe_fmt_pct(stats_row.get(f"Jensen Alpha ({bname})"))])
            rows.append([f"Treynor ({bname})", _safe_fmt_float(stats_row.get(f"Treynor ({bname})"))])
        rows.append(["CVaR 99% (Daily)", _safe_fmt_pct(stats_row.get("CVaR 99% (Daily)"))])
        rows.append(["Max Drawdown", _safe_fmt_pct(stats_row.get("Max Drawdown"))])
        t = Table([["Metric", "Value"], *rows], colWidths=[110 * mm, 45 * mm])
        t.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), DARK_TEAL),
                    ("TEXTCOLOR", (0, 0), (-1, 0), WHITE),
                    ("FONTNAME", (0, 0), (-1, 0), TITLE_STYLE.fontName),
                    ("FONTNAME", (0, 1), (-1, -1), BODY_STYLE.fontName),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [WHITE, LIGHT_GRAY]),
                    ("GRID", (0, 0), (-1, -1), 0.3, MID_GRAY),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]
            )
        )
        stats_content = t
    else:
        stats_content = Paragraph("Statistics not available.", BODY_STYLE)

    story.append(KeepTogether([
        Paragraph("Summary Statistics", SECTION_STYLE),
        HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.5),
        Spacer(1, 4 * mm),
        stats_content,
    ]))
    story.append(Spacer(1, 8 * mm))

    story.append(KeepTogether([
        Paragraph("Performance", SECTION_STYLE),
        HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.5),
        Spacer(1, 4 * mm),
        growth_chart(returns.dropna(), str(portfolio_name), benchmark_returns=benchmark_returns),
    ]))
    story.append(Spacer(1, 6 * mm))
    story.append(drawdown_chart(returns.dropna(), str(portfolio_name)))

    story.append(Spacer(1, 8 * mm))
    tr_tbl = trailing_returns_pdf_table(trailing_returns) if trailing_returns is not None else None
    tr_content = tr_tbl if tr_tbl is not None else Paragraph("Not enough data for trailing period returns.", BODY_STYLE)
    story.append(KeepTogether([
        Paragraph("Trailing Period Returns", SUBTITLE_STYLE),
        tr_content,
    ]))

    story.append(Spacer(1, 8 * mm))

    tbl = top_holdings_table(holdings, prices, n=15)
    holdings_content = tbl if tbl is not None else Paragraph("No holdings available.", BODY_STYLE)
    story.append(KeepTogether([
        Paragraph("Holdings & Allocation", SECTION_STYLE),
        HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.5),
        Spacer(1, 4 * mm),
        holdings_content,
    ]))

    story.append(Spacer(1, 10 * mm))
    sec = sector_pie_image(holdings, prices, sector_map=sector_map)
    sector_content = sec if sec is not None else Paragraph("Sector map not available.", BODY_STYLE)
    story.append(KeepTogether([
        Paragraph("Sector Exposure", SECTION_STYLE),
        HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.5),
        Spacer(1, 4 * mm),
        sector_content,
    ]))

    story.append(Spacer(1, 10 * mm))
    cap_img = cap_split_pie_image(cap_split) if cap_split is not None else None
    cap_content = cap_img if cap_img is not None else Paragraph("Market cap data not available.", BODY_STYLE)
    story.append(KeepTogether([
        Paragraph("Market Cap Allocation", SECTION_STYLE),
        HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.5),
        Spacer(1, 4 * mm),
        cap_content,
    ]))

    story.append(Spacer(1, 8 * mm))
    rr_df = stock_risk_return(
        attribution if attribution is not None else pd.Series(dtype=float),
        risk_contrib if risk_contrib is not None else pd.Series(dtype=float),
    )
    rr_fig = risk_return_scatter_figure(rr_df) if not rr_df.empty else None
    rr_content = (
        make_chart_image(rr_fig, width_mm=150, height_mm=115)
        if rr_fig is not None
        else Paragraph("Not enough attribution/risk data to plot risk vs. return contribution.", BODY_STYLE)
    )
    story.append(KeepTogether([
        Paragraph("Risk vs. Return Contribution by Holding", SECTION_STYLE),
        HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.5),
        Spacer(1, 4 * mm),
        rr_content,
    ]))

    story.append(Spacer(1, 8 * mm))

    port_vol_for_rc = annualized_volatility(returns.dropna()) if returns is not None else np.nan
    rc_tbl = (
        risk_contribution_table(risk_contrib, port_vol_for_rc)
        if risk_contrib is not None and not risk_contrib.empty
        else None
    )
    rc_content = rc_tbl if rc_tbl is not None else Paragraph("Risk contribution not available.", BODY_STYLE)
    story.append(KeepTogether([
        Paragraph("Risk & Attribution", SECTION_STYLE),
        HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.5),
        Spacer(1, 4 * mm),
        rc_content,
    ]))

    story.append(Spacer(1, 8 * mm))
    at_tbl = contribution_table(attribution) if attribution is not None else None
    at_content = at_tbl if at_tbl is not None else Paragraph("Performance attribution not available.", BODY_STYLE)
    story.append(KeepTogether([
        Paragraph("Contribution to Return (Since Inception)", SUBTITLE_STYLE),
        at_content,
    ]))

    story.append(Spacer(1, 8 * mm))
    hist_img = returns_histogram_image(returns_hist, "Histogram of daily returns") if returns_hist is not None else None
    hist_content = hist_img if hist_img is not None else Paragraph("Return distribution not available.", BODY_STYLE)
    hist_stats_tbl = returns_stats_table(returns_hist) if returns_hist is not None else None
    return_dist_group = [
        Paragraph("Return Distribution", SECTION_STYLE),
        HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.5),
        Spacer(1, 4 * mm),
        hist_content,
    ]
    if hist_stats_tbl is not None:
        return_dist_group += [Spacer(1, 4 * mm), hist_stats_tbl]
    story.append(KeepTogether(return_dist_group))

    story.append(Spacer(1, 8 * mm))
    corr_img = corr_heatmap_image(stock_corr, "Correlation heatmap (stocks)") if stock_corr is not None else None
    corr_content = corr_img if corr_img is not None else Paragraph("Correlation heatmap not available.", BODY_STYLE)
    n500_returns_for_corr = (benchmark_returns or {}).get("Nifty 500")
    corr_summary = correlation_summary(stock_corr, returns, n500_returns_for_corr)
    corr_summary_text = (
        f"Average pairwise correlation among this portfolio's holdings: {_safe_fmt_float(corr_summary['avg_pairwise'])}. "
        f"Portfolio correlation vs. Nifty 500: {_safe_fmt_float(corr_summary['vs_benchmark'])}."
    )
    story.append(KeepTogether([
        Paragraph("Stock Correlation", SECTION_STYLE),
        HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.5),
        Spacer(1, 4 * mm),
        corr_content,
        Spacer(1, 3 * mm),
        Paragraph(corr_summary_text, BODY_STYLE),
    ]))

    story.append(Spacer(1, 8 * mm))
    story.append(KeepTogether([
        Paragraph("Disclaimer", SECTION_STYLE),
        HRFlowable(width="100%", color=ACCENT_TEAL, thickness=0.5),
        Spacer(1, 4 * mm),
        Paragraph(
            "Equity investments are subject to market risks. Past performance is not indicative of future returns. "
            "This report is for information purposes only and should not be construed as investment advice.",
            BODY_STYLE,
        ),
    ]))

    doc.build(story, onFirstPage=_page_header, onLaterPages=_page_header)
    return buf.getvalue()


def generate_portfolio_pptx_bytes(
    *,
    portfolio_name: str,
    report_date: date | datetime,
    returns: pd.Series,
    holdings: pd.DataFrame,
    prices: pd.DataFrame,
    stats_row: Optional[pd.Series] = None,
    benchmark_returns: Optional[dict[str, pd.Series]] = None,
    sector_map: Optional[dict[str, str]] = None,
    risk_contrib: Optional[pd.Series] = None,
    attribution: Optional[pd.Series] = None,
    stock_corr: Optional[pd.DataFrame] = None,
    returns_hist: Optional[pd.Series] = None,
    cash_balance: Optional[float] = None,
    portfolio_value: Optional[float] = None,
    cap_split: Optional[pd.Series] = None,
    trailing_returns: Optional[pd.DataFrame] = None,
) -> bytes:
    """
    Same data as generate_portfolio_pdf_bytes (same parameters, on purpose
    — one call site can feed both), rendered as an editable PowerPoint
    instead. Text and tables are genuine native pptx elements — the client
    team can retype numbers, restyle a table, or drop in their own
    commentary directly in PowerPoint. Charts are embedded as images, the
    same convention every BI tool's "export to PowerPoint" uses — a fully
    live-editable native chart tied to this exact styling isn't something
    python-pptx can reproduce from a matplotlib figure, so this doesn't
    pretend otherwise.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    BRAND_INDIGO = RGBColor(0x29, 0x00, 0x84)
    BRAND_GOLD = RGBColor(0xF6, 0xBF, 0x02)
    BRAND_GREEN = RGBColor(0x1D, 0x99, 0x51)
    BRAND_RED = RGBColor(0xAA, 0x14, 0x37)
    WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY_RGB = RGBColor(0xF2, 0xF3, 0xF4)
    DARK_TEXT = RGBColor(0x1B, 0x1B, 0x1B)
    MID_GRAY_RGB = RGBColor(0x71, 0x7D, 0x7E)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    _assets = Path(__file__).resolve().parent / "assets"
    logo_file = next(
        (p for p in [
            _assets / "vika_logo.png", _assets / "vika_logo.jpg",
            _assets / "Vika_Logo.png", _assets / "Vika_Logo.jpg",
        ] if p.exists()),
        None,
    )

    def _add_slide():
        slide = prs.slides.add_slide(blank_layout)
        if logo_file is not None:
            try:
                slide.shapes.add_picture(str(logo_file), Inches(11.6), Inches(0.25), height=Inches(0.7))
            except Exception:
                pass
        return slide

    def _add_title(slide, text, top=Inches(0.3)):
        tb = slide.shapes.add_textbox(Inches(0.5), top, Inches(10.5), Inches(0.6))
        tf = tb.text_frame
        tf.text = text
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = BRAND_INDIGO
        return tb

    def _add_body_text(slide, text, top, height=Inches(2.5), size=14):
        tb = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = text
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(size)
        run.font.color.rgb = DARK_TEXT
        return tb

    def _style_table(tbl, n_rows, n_cols):
        for c in range(n_cols):
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND_INDIGO
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.color.rgb = WHITE_RGB
                p.font.size = Pt(12)
        for r in range(1, n_rows):
            for c in range(n_cols):
                cell = tbl.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY_RGB if r % 2 == 0 else WHITE_RGB
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = DARK_TEXT

    def _add_table(slide, rows, top, col_widths=None, left=Inches(0.5), width=Inches(12.3)):
        n_rows, n_cols = len(rows), len(rows[0])
        height = Inches(0.4 * n_rows)
        gframe = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        tbl = gframe.table
        if col_widths:
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = w
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                tbl.cell(r, c).text = str(val)
        _style_table(tbl, n_rows, n_cols)
        return gframe

    def _add_image_fig(slide, fig, top, left=None, width=Inches(9)):
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
        buf.seek(0)
        plt.close(fig)
        if left is None:
            left = Inches((13.333 - width.inches) / 2)
        slide.shapes.add_picture(buf, left, top, width=width)

    # ---- Slide 1: Title ----
    slide = _add_slide()
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(1.2))
    tf = tb.text_frame
    tf.text = "Portfolio Report"
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = BRAND_INDIGO
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb2 = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(11.3), Inches(0.6))
    tf2 = tb2.text_frame
    tf2.text = str(portfolio_name)
    run2 = tf2.paragraphs[0].runs[0]
    run2.font.size = Pt(20)
    run2.font.color.rgb = MID_GRAY_RGB
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb3 = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(11.3), Inches(0.5))
    tf3 = tb3.text_frame
    tf3.text = f"Report date: {pd.to_datetime(report_date).date().isoformat()}"
    run3 = tf3.paragraphs[0].runs[0]
    run3.font.size = Pt(14)
    run3.font.color.rgb = MID_GRAY_RGB
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ---- Slide 2: Executive Summary + Summary Statistics ----
    slide = _add_slide()
    _add_title(slide, "Executive Summary")
    _add_body_text(
        slide,
        summary_text(portfolio_name=str(portfolio_name), stats_row=stats_row, returns=returns),
        top=Inches(1.1), height=Inches(1.3), size=13,
    )
    if stats_row is not None:
        rows = [["Metric", "Value"]]
        if portfolio_value is not None:
            rows.append(["Current Portfolio Value", fmt_inr(portfolio_value)])
        if cash_balance is not None:
            rows.append(["Cash Balance", fmt_inr(cash_balance)])
        rows += [
            ["Period Return (TWR)",     _safe_fmt_pct(stats_row.get("Period Return (TWR)"))],
            ["Period Volatility (TWR)", _safe_fmt_pct(stats_row.get("Period Volatility (TWR)"))],
            ["Sharpe",  _safe_fmt_float(stats_row.get("Sharpe"))],
            ["Sortino", _safe_fmt_float(stats_row.get("Sortino"))],
        ]
        bench_names = sorted({
            k[len("Beta ("):-1] for k in stats_row.index if k.startswith("Beta (") and k.endswith(")")
        })
        for bname in bench_names:
            rows.append([f"Beta ({bname})", _safe_fmt_float(stats_row.get(f"Beta ({bname})"))])
            rows.append([f"Jensen Alpha ({bname}) (TWR trailing period)", _safe_fmt_pct(stats_row.get(f"Jensen Alpha ({bname})"))])
            rows.append([f"Treynor ({bname})", _safe_fmt_float(stats_row.get(f"Treynor ({bname})"))])
        rows.append(["CVaR 99% (Daily)", _safe_fmt_pct(stats_row.get("CVaR 99% (Daily)"))])
        rows.append(["Max Drawdown", _safe_fmt_pct(stats_row.get("Max Drawdown"))])
        _add_table(slide, rows, top=Inches(2.6), col_widths=[Inches(8.5), Inches(3.8)])

    # ---- Slide 3: Performance chart ----
    slide = _add_slide()
    _add_title(slide, "Performance")
    fig = growth_chart_figure(returns.dropna(), str(portfolio_name), benchmark_returns=benchmark_returns)
    if fig is not None:
        _add_image_fig(slide, fig, top=Inches(1.1), width=Inches(11.5))

    # ---- Slide 4: Trailing Period Returns ----
    if trailing_returns is not None and not trailing_returns.empty:
        slide = _add_slide()
        _add_title(slide, "Trailing Period Returns")
        header = [""] + list(trailing_returns.columns)
        rows = [header]
        for label, row in trailing_returns.iterrows():
            rows.append([label] + ["" if pd.isna(v) else f"{v*100:.2f}%" for v in row])
        _add_table(slide, rows, top=Inches(1.3))

    # ---- Slide 5: Holdings & Allocation ----
    slide = _add_slide()
    _add_title(slide, "Holdings & Allocation")
    tickers_h = [t for t in holdings.columns if t in prices.columns]
    if tickers_h:
        latest_h = holdings.iloc[-1][tickers_h]
        latest_p = prices.loc[prices.index[-1], tickers_h]
        mv = (latest_h * latest_p).sort_values(ascending=False)
        total_mv = mv.sum()
        rows = [["Ticker", "Market Value (INR)", "Weight"]]
        for ticker, val in mv.head(15).items():
            weight_str = f"{val/total_mv*100:.1f}%" if total_mv and abs(total_mv) > 1e-9 else "N/A"
            rows.append([ticker.replace(".NS", "").replace(".BO", ""), fmt_inr(val), weight_str])
        _add_table(slide, rows, top=Inches(1.1))

    # ---- Slide 6: Sector Exposure ----
    if sector_map is not None:
        exp, counts = sector_exposure_weights_and_counts(holdings, prices, sector_map or {})
        if not exp.empty:
            slide = _add_slide()
            _add_title(slide, "Sector Exposure")
            fig, colors_hex_pptx = sector_pie_figure_clean(exp)
            _add_image_fig(slide, fig, top=Inches(1.2), left=Inches(0.6), width=Inches(5.5))

            n_rows = len(exp) + 1
            tbl_top, tbl_left = Inches(1.2), Inches(6.6)
            tbl_width, tbl_height = Inches(6), Inches(0.35 * n_rows)
            gframe = slide.shapes.add_table(n_rows, 4, tbl_left, tbl_top, tbl_width, tbl_height)
            tbl = gframe.table
            tbl.columns[0].width = Inches(0.4)
            tbl.columns[1].width = Inches(3.2)
            tbl.columns[2].width = Inches(1.2)
            tbl.columns[3].width = Inches(1.2)
            headers = ["", "Sector", "%", "Stocks"]
            for c, h in enumerate(headers):
                cell = tbl.cell(0, c)
                cell.text = h
                cell.fill.solid()
                cell.fill.fore_color.rgb = BRAND_INDIGO
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.color.rgb = WHITE_RGB
                    p.font.size = Pt(11)
            for r, (sector, weight) in enumerate(exp.items(), start=1):
                swatch_cell = tbl.cell(r, 0)
                swatch_cell.fill.solid()
                hexcolor = colors_hex_pptx[r - 1].lstrip("#")
                swatch_cell.fill.fore_color.rgb = RGBColor(int(hexcolor[0:2], 16), int(hexcolor[2:4], 16), int(hexcolor[4:6], 16))
                vals = [sector, f"{weight*100:.1f}%", str(int(counts.get(sector, 0)))]
                for c, v in enumerate(vals, start=1):
                    cell = tbl.cell(r, c)
                    cell.text = str(v)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY_RGB if r % 2 == 0 else WHITE_RGB
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(10)
                        p.font.color.rgb = DARK_TEXT

    # ---- Slide 7: Market Cap Allocation ----
    cap_split_clean = normalize_cap_split_for_report(cap_split)
    if cap_split_clean is not None and not cap_split_clean.empty:
        slide = _add_slide()
        _add_title(slide, "Market Cap Allocation")
        color_map = {"Large Cap": "#290084", "Mid Cap": "#F6BF02", "Small Cap": "#1D9951", "ETF": "#717D7E", "REIT": "#AA1437", "Unknown": "#B0B0B0"}
        colors_ordered = [color_map.get(k, "#B0B0B0") for k in cap_split_clean.index]
        fig, ax = plt.subplots(figsize=(6, 6))
        _, _, autotexts = ax.pie(
            cap_split_clean.values, labels=cap_split_clean.index, autopct="%1.1f%%", startangle=90,
            colors=colors_ordered, textprops={"fontsize": 9},
        )
        for txt, color in zip(autotexts, colors_ordered):
            txt.set_color(_contrast_text_color(color))
        fig.tight_layout()
        _add_image_fig(slide, fig, top=Inches(1.1), width=Inches(6.5))

    # ---- Slide 8: Risk vs. Return Contribution ----
    rr_df = stock_risk_return(
        attribution if attribution is not None else pd.Series(dtype=float),
        risk_contrib if risk_contrib is not None else pd.Series(dtype=float),
    )
    if not rr_df.empty:
        slide = _add_slide()
        _add_title(slide, "Risk vs. Return Contribution by Holding")
        fig = risk_return_scatter_figure(rr_df)
        if fig is not None:
            _add_image_fig(slide, fig, top=Inches(1.1), width=Inches(9))

    # ---- Slide 9: Risk & Attribution tables ----
    if (risk_contrib is not None and not risk_contrib.empty) or (attribution is not None and not attribution.empty):
        slide = _add_slide()
        _add_title(slide, "Risk & Attribution")
        top_cursor = Inches(1.1)
        if risk_contrib is not None and not risk_contrib.empty:
            port_vol = annualized_volatility(returns.dropna()) if returns is not None else np.nan
            s = pd.to_numeric(risk_contrib, errors="coerce").dropna().sort_values(ascending=False).head(8)
            rows = [["Ticker", "Contribution to Volatility", "% of Total Risk"]]
            for ticker, share in s.items():
                vol_c = f"{share*port_vol*100:+.2f}%" if pd.notna(port_vol) else "N/A"
                rows.append([str(ticker).replace(".NS", "").replace(".BO", ""), vol_c, f"{share*100:+.2f}%"])
            _add_table(slide, rows, top=top_cursor, width=Inches(6), left=Inches(0.5))
        if attribution is not None and not attribution.empty:
            s = pd.to_numeric(attribution, errors="coerce").dropna().sort_values(ascending=False).head(8)
            total = s.sum()
            rows = [["Ticker", "Contribution to Return", "% of Total Return"]]
            for ticker, val in s.items():
                pct_total = f"{val/total*100:+.2f}%" if total and abs(total) > 1e-9 else "N/A"
                rows.append([str(ticker).replace(".NS", "").replace(".BO", ""), f"{val*100:+.2f}%", pct_total])
            _add_table(slide, rows, top=top_cursor, width=Inches(6), left=Inches(6.8))

    # ---- Slide 10: Return Distribution ----
    if returns_hist is not None:
        r = pd.to_numeric(returns_hist, errors="coerce").dropna()
        if len(r) >= 10:
            slide = _add_slide()
            _add_title(slide, "Return Distribution")
            fig = returns_histogram_figure(r)
            if fig is not None:
                _add_image_fig(slide, fig, top=Inches(1.1), width=Inches(8))
            rows = [
                ["Statistic", "Value"],
                ["Mean (daily)", f"{float(r.mean())*100:+.3f}%"],
                ["Std. Dev. (daily)", f"{float(r.std(ddof=1))*100:.3f}%"],
                ["Skewness", f"{float(r.skew()):+.2f}"],
                ["Kurtosis (excess)", f"{float(r.kurt()):+.2f}"],
                ["Min (daily)", f"{float(r.min())*100:+.2f}%"],
                ["Max (daily)", f"{float(r.max())*100:+.2f}%"],
            ]
            _add_table(slide, rows, top=Inches(1.1), left=Inches(9.3), width=Inches(3.5), col_widths=[Inches(2), Inches(1.5)])

    # ---- Slide 11: Stock Correlation ----
    if stock_corr is not None and not stock_corr.empty:
        slide = _add_slide()
        _add_title(slide, "Stock Correlation")
        fig = corr_heatmap_figure(stock_corr)
        if fig is not None:
            _add_image_fig(slide, fig, top=Inches(1.0), width=Inches(7.5))
        n500_returns = (benchmark_returns or {}).get("Nifty 500")
        csum = correlation_summary(stock_corr, returns, n500_returns)
        _add_body_text(
            slide,
            f"Average pairwise correlation among this portfolio's holdings: {_safe_fmt_float(csum['avg_pairwise'])}. "
            f"Portfolio correlation vs. Nifty 500: {_safe_fmt_float(csum['vs_benchmark'])}.",
            top=Inches(6.7), height=Inches(0.5), size=12,
        )

    # ---- Slide 12: Disclaimer ----
    slide = _add_slide()
    _add_title(slide, "Disclaimer")
    _add_body_text(
        slide,
        "Equity investments are subject to market risks. Past performance is not indicative of future returns. "
        "This report is for information purposes only and should not be construed as investment advice.",
        top=Inches(1.2), height=Inches(1.0), size=13,
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ============================== SECTION: app (Streamlit UI) ==============================

st.set_page_config(page_title="Portfolio Stats", layout="wide")

# Visible build marker — if this string isn't showing up in the running app
# after a redeploy, the deploy didn't actually pick up the new file. Bump
# this on every future change so it's always possible to confirm at a
# glance whether the live app matches what was last sent.
APP_BUILD_TAG = "build-2026-08-19-sector-pie-legend-table-16"
st.caption(f"⚙ {APP_BUILD_TAG}")


def _df_to_csv_bytes(df: pd.DataFrame) -> bytes:
    return df.to_csv(index=True).encode("utf-8")


def _fmt_pct(x) -> str:
    try:
        v = float(x)
    except Exception:
        return "N/A"
    if pd.isna(v) or v in (float("inf"), float("-inf")):
        return "N/A"
    return f"{v*100:,.2f}%"


def _fmt_float(x) -> str:
    try:
        v = float(x)
    except Exception:
        return "N/A"
    if pd.isna(v) or v in (float("inf"), float("-inf")):
        return "N/A"
    return f"{v:,.2f}"


def _fmt_inr(x) -> str:
    return format_inr(x) or "N/A"


def _fmt_pct_signed(x) -> str:
    try:
        v = float(x)
    except Exception:
        return "N/A"
    if pd.isna(v) or v in (float("inf"), float("-inf")):
        return "N/A"
    return f"{v*100:+,.2f}%"


@st.cache_data(show_spinner=False, ttl=60 * 60)
def _download_prices_cached(
    tickers_ns: tuple[str, ...], start_iso: str, include_benchmark: bool, benchmark_tickers: tuple[str, ...]
):
    import yfinance as yf

    start_dt = pd.to_datetime(start_iso).to_pydatetime()
    return download_prices(
        yf,
        tickers_ns=tickers_ns,
        start=start_dt,
        include_benchmark=include_benchmark,
        benchmark_tickers=benchmark_tickers,
    )


@st.cache_data(show_spinner=False, ttl=24 * 60 * 60)
def _fetch_market_caps_cached(tickers_ns: tuple[str, ...]) -> pd.Series:
    """Best-effort market cap (INR) per ticker via yfinance, for the market-cap
    allocation pie chart. Cached for a day since market cap barely moves
    intraday and this is one Ticker() call per holding."""
    import yfinance as yf

    return fetch_market_caps(yf, tickers_ns)




# Show logo — support both png and jpg
import os as _os
_logo_path = next(
    (p for p in ["assets/vika_logo.png", "assets/vika_logo.jpg", "assets/Vika_Logo.jpg", "assets/Vika_Logo.png"]
     if _os.path.exists(p)),
    None,
)
if _logo_path:
    try:
        st.image(_logo_path, width=200)
    except Exception:
        pass

st.title("Portfolio Stats")
st.caption("Upload a transactions Excel file → view portfolio analytics → download reports per client.")


def _bar_chart_from_series(
    series: pd.Series, *, title: str, value_label: str = "Value", top_n: int = 12, as_percent: bool = True
):
    import altair as alt

    s = pd.to_numeric(series, errors="coerce").dropna()
    if s.empty:
        st.info("Not enough data to chart.")
        return
    s = s.sort_values(ascending=False).head(top_n)
    display_vals = s.values * 100 if as_percent else s.values
    dfc = pd.DataFrame({"Label": s.index.astype(str), value_label: display_vals})
    fmt = ",.2f" if as_percent else ",.6f"
    axis_title = f"{value_label} (%)" if as_percent else value_label
    chart = (
        alt.Chart(dfc)
        .mark_bar()
        .encode(
            y=alt.Y("Label:N", sort="-x", title=None),
            x=alt.X(f"{value_label}:Q", title=axis_title),
            tooltip=["Label:N", alt.Tooltip(f"{value_label}:Q", format=fmt, title=axis_title)],
        )
        .properties(height=min(28 * len(dfc), 360), title=title)
    )
    st.altair_chart(chart, use_container_width=True)


def _heatmap(df: pd.DataFrame, *, title: str):
    import matplotlib.pyplot as plt

    if df.empty:
        st.info("Not enough data.")
        return

    cmap = get_corr_cmap()
    cmap = cmap.copy() if hasattr(cmap, "copy") else cmap
    cmap.set_bad(color="white")

    # Only the lower triangle (plus diagonal) — a correlation matrix is
    # symmetric, so the upper triangle is a mirror image that adds nothing.
    values = df.values.astype(float).copy()
    if values.shape[0] == values.shape[1]:
        upper = np.triu_indices(values.shape[0], k=1)
        values[upper] = np.nan

    fig, ax = plt.subplots(figsize=(10, 6))
    im = ax.imshow(values, aspect="auto", cmap=cmap, vmin=-1, vmax=1)
    ax.set_xticks(np.arange(len(df.columns)))
    ax.set_yticks(np.arange(len(df.index)))
    ax.set_xticklabels(df.columns, rotation=45, ha="right")
    ax.set_yticklabels(df.index)
    for i in range(len(df.index)):
        for j in range(len(df.columns)):
            v = values[i, j]
            if pd.isna(v):
                continue
            # Luminance-based contrast: the navy/gold colormap isn't symmetric
            # in brightness (navy darkens faster than gold), so a fixed |v|
            # cutoff picks the wrong text color near +1; check actual cell
            # luminance instead.
            r, g, b, _ = cmap((float(v) + 1) / 2)
            luminance = 0.299 * r + 0.587 * g + 0.114 * b
            text_color = "white" if luminance < 0.55 else "black"
            ax.text(j, i, f"{v:.2f}", ha="center", va="center", fontsize=8, color=text_color)
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    ax.set_title(title)
    fig.tight_layout()
    st.pyplot(fig, use_container_width=True)
    plt.close(fig)

with st.sidebar:
    st.header("Upload")
    uploaded = st.file_uploader("Transactions Excel (.xlsx)", type=["xlsx"])
    password = st.text_input("Password (only if file is encrypted)", type="password")
    include_benchmark = st.checkbox("Include benchmarks (Nifty 50 & Nifty 500)", value=True)
    st.divider()
    st.header("Assumptions")
    rf_annual = st.number_input("Risk-free rate (annual, %)", min_value=0.0, value=0.0, step=0.25) / 100.0
    st.caption(
        "Defaults to 0% — Sharpe, Sortino, and Treynor all subtract this from the return, so at 0% those "
        "ratios show raw return/risk with no risk-free adjustment, which reads high/unusual next to "
        "reports that use a real rate. Set this to your actual reference rate (e.g. a 91-day T-bill or "
        "repo-linked yield) to get comparable ratios."
    )

if not uploaded:
    st.info("Upload the transactions Excel file to begin.")
    st.stop()

try:
    excel_bytes = uploaded.getvalue()
    sheets = load_transactions_excel(excel_bytes, password=password or None)
    # prepare_transactions reads all three sheets — "Initial Holdings" (seeded as
    # opening BUY transactions at each portfolio's StartDate), "Transactions"
    # (the trade log), and "Cashflows" (deposits/withdrawals, kept separate).
    df, cashflows_df = prepare_transactions(sheets)
    # Any SELL that exceeds recorded holdings (no matching BUY in Initial
    # Holdings or Transactions) gets an inferred opening position backfilled
    # automatically, so it's correctly reflected everywhere downstream
    # instead of silently clamped to zero — see auto_backfill_oversold_positions.
    df, backfill_summary = auto_backfill_oversold_positions(df)
except TransactionsFormatError as e:
    st.error(str(e))
    st.stop()
except Exception as e:
    st.error(f"Failed to read the uploaded file: {e}")
    st.stop()

positions = current_positions(df)
# ALL ever-transacted tickers, not just current_positions()'s currently-held
# subset — a fully-exited position (bought and later fully sold) still
# needs price history for attribution/P&L over the period it was held, and
# still shows up as a column in daily_positions (which tracks every ticker
# ever transacted, defaulting to 0 once sold out). Downloading only the
# currently-held subset meant a fully-exited ticker would show up in
# "missing from downloaded prices" despite never actually being attempted
# — which is also why its failure_reasons entry was always empty.
tickers = tuple(sorted(df["ticker_yf"].unique().tolist()))
start_iso = df["date"].min().date().isoformat()

with st.spinner("Downloading market data (yfinance)…"):
    dl = _download_prices_cached(tickers, start_iso, include_benchmark, tuple(BENCHMARKS.values()))

prices = dl.prices.sort_index()
if prices.empty:
    st.error("No prices could be downloaded. Check ticker mappings and internet access on the host.")
    st.stop()

# One single position/cash/equity walk — build_daily_ledger both produces
# the TWR ledger AND the per-day share-quantity history (daily_positions)
# used everywhere else (attribution, risk contribution, sector exposure,
# market-cap allocation, the holdings table). There's no longer a second,
# independently-computed position tracker to disagree with it.
ledgers, daily_positions, oversold_warnings = build_daily_ledger(df, cashflows_df, prices)
missing_tickers = missing_price_tickers(daily_positions, prices)
# TWR returns come from a full daily ledger (equity + cash combined), not
# stock prices alone — deposits/withdrawals are excluded from the return via
# the External Flow adjustment in compute_ledger_twr, so they still never
# show up as investment gain/loss, but idle cash and cash timing are now
# properly reflected in the portfolio's value and TWR.
twr_series = {}
for _portfolio, _ledger in ledgers.items():
    twr_series[_portfolio] = compute_ledger_twr(_ledger)
twr_df = pd.DataFrame(twr_series)
twr_df.columns = twr_df.columns.astype(str)
portfolios_all = sorted(str(p) for p in daily_positions.keys())

benchmarks = (
    {name: prices[ticker] for name, ticker in BENCHMARKS.items() if ticker in prices.columns}
    if include_benchmark
    else {}
)

overview_tab, client_tab = st.tabs(["Overview", "Client report"])

with overview_tab:
    st.subheader("Key metrics")
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Clients", f"{len(portfolios_all)}")
    c2.metric("Tickers", f"{len(tickers)}")
    c3.metric("From", f"{df['date'].min().date().isoformat()}")
    c4.metric("To", f"{df['date'].max().date().isoformat()}")

    st.divider()
    st.subheader("Portfolio statistics")

    # Compute metrics here so rf_annual slider changes are reflected immediately
    metrics_rows = []
    for p in portfolios_all:
        twr_col = twr_df[p].dropna() if p in twr_df.columns else pd.Series(dtype=float)
        m = portfolio_metrics(twr_col, benchmarks=benchmarks, rf_annual=rf_annual)
        m.name = str(p)
        metrics_rows.append(m)
    stats = pd.DataFrame(metrics_rows)

    stats_disp = stats.copy()
    for col in stats_disp.columns:
        s = pd.to_numeric(stats_disp[col], errors="coerce")
        stats_disp[col] = s.mask(~pd.Series(np.isfinite(s), index=s.index), np.nan)

    # Show a clean, formatted table instead of lots of "None". Values well
    # outside what's plausible for a diversified equity portfolio get a ⚠
    # marker rather than being shown as clean-looking numbers — those are
    # almost always downstream of a data issue flagged in Data Quality
    # below (missing prices, negative cash, thin benchmark data), not a
    # real result.
    def _fmt_col_pct(c: pd.Series, lo: float = None, hi: float = None) -> pd.Series:
        def _f(v):
            if pd.isna(v):
                return "N/A"
            s = f"{v*100:,.2f}%"
            if (lo is not None and v < lo) or (hi is not None and v > hi):
                return f"⚠ {s}"
            return s
        return c.apply(_f)

    def _fmt_col_float(c: pd.Series, lo: float = None, hi: float = None) -> pd.Series:
        def _f(v):
            if pd.isna(v):
                return "N/A"
            s = f"{v:,.2f}"
            if (lo is not None and v < lo) or (hi is not None and v > hi):
                return f"⚠ {s}"
            return s
        return c.apply(_f)

    stats_fmt = pd.DataFrame(index=stats_disp.index.astype(str))
    if "Period Return (TWR)" in stats_disp.columns:
        stats_fmt["Period Return (TWR)"] = _fmt_col_pct(stats_disp["Period Return (TWR)"], lo=-0.95, hi=5.0)
    if "Period Volatility (TWR)" in stats_disp.columns:
        stats_fmt["Period Volatility (TWR)"] = _fmt_col_pct(stats_disp["Period Volatility (TWR)"], lo=0, hi=1.5)
    if "Sharpe" in stats_disp.columns:
        stats_fmt["Sharpe"] = _fmt_col_float(stats_disp["Sharpe"], lo=-10, hi=10)
    if "Sortino" in stats_disp.columns:
        stats_fmt["Sortino"] = _fmt_col_float(stats_disp["Sortino"], lo=-10, hi=10)
    # Beta/Jensen Alpha/Treynor exist once per benchmark (e.g. "Beta (Nifty 50)",
    # "Beta (Nifty 500)") — pick them up dynamically rather than hardcoding one.
    for col in stats_disp.columns:
        if col.startswith("Beta (") or col.startswith("Treynor ("):
            stats_fmt[col] = _fmt_col_float(stats_disp[col], lo=-5, hi=5)
        elif col.startswith("Jensen Alpha ("):
            stats_fmt[f"{col} (TWR trailing period)"] = _fmt_col_pct(stats_disp[col], lo=-2.0, hi=2.0)
    if "CVaR 99% (Daily)" in stats_disp.columns:
        stats_fmt["CVaR 99% (Daily)"] = _fmt_col_pct(stats_disp["CVaR 99% (Daily)"])
    if "Max Drawdown" in stats_disp.columns:
        stats_fmt["Max Drawdown"] = _fmt_col_pct(stats_disp["Max Drawdown"])

    st.dataframe(stats_fmt, use_container_width=True)
    if stats_fmt.apply(lambda col: col.astype(str).str.startswith("⚠")).to_numpy().any():
        st.caption(
            "⚠ = outside a plausible range for a diversified equity portfolio. Check Data Quality below "
            "for that portfolio (missing prices, negative cash, or a stale benchmark are the usual causes) "
            "before using these numbers."
        )

    # ── Realized vs Unrealized P&L ──────────────────────────────────────────
    st.subheader("Realized vs Unrealized P&L (₹)")
    with st.spinner("Computing P&L…"):
        pnl_df = compute_pnl(df, prices)

    if pnl_df.empty:
        st.info("No P&L data available.")
    else:
        import matplotlib.pyplot as plt
        import matplotlib.ticker as mticker

        portfolios_pnl = pnl_df.index.tolist()
        x = np.arange(len(portfolios_pnl))
        width = 0.32

        fig_pnl, ax_pnl = plt.subplots(figsize=(max(7, len(portfolios_pnl) * 1.6), 4.5))
        bars_r = ax_pnl.bar(x - width / 2, pnl_df["realized_pnl"], width,
                            label="Realized P&L",
                            color=["#27ae60" if v >= 0 else "#e74c3c" for v in pnl_df["realized_pnl"]])
        bars_u = ax_pnl.bar(x + width / 2, pnl_df["unrealized_pnl"], width,
                            label="Unrealized P&L",
                            color=["#2980b9" if v >= 0 else "#e67e22" for v in pnl_df["unrealized_pnl"]])

        ax_pnl.set_xticks(x)
        ax_pnl.set_xticklabels(portfolios_pnl, rotation=30, ha="right", fontsize=9)
        ax_pnl.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"₹{v:,.0f}"))
        ax_pnl.axhline(0, color="black", linewidth=0.8, linestyle="--")
        ax_pnl.set_title("Realized vs Unrealized P&L per Portfolio", fontsize=11)
        ax_pnl.legend(fontsize=9)
        ax_pnl.grid(False)

        # Label bars with values
        for bar in list(bars_r) + list(bars_u):
            h = bar.get_height()
            if abs(h) > 0:
                ax_pnl.text(
                    bar.get_x() + bar.get_width() / 2,
                    h + (max(pnl_df[["realized_pnl", "unrealized_pnl"]].abs().max()) * 0.01),
                    f"₹{h:,.0f}",
                    ha="center", va="bottom", fontsize=7, rotation=45,
                )

        fig_pnl.tight_layout()
        st.pyplot(fig_pnl, use_container_width=True)
        plt.close(fig_pnl)

        # Also show as table
        pnl_display = pnl_df.copy()
        for col in pnl_display.columns:
            pnl_display[col] = pnl_display[col].apply(lambda v: f"₹{v:,.2f}")
        pnl_display.columns = ["Realized P&L", "Unrealized P&L", "Total P&L"]
        st.dataframe(pnl_display, use_container_width=True)

    # ── Cash flows ───────────────────────────────────────────────────────────
    st.subheader("Cash flows (capital contributed)")
    if cashflows_df.empty:
        st.info("No cashflow data found.")
    else:
        summary = cashflow_summary(cashflows_df)
        summary_disp = summary.to_frame("Net cash contributed (₹)")
        summary_disp["Net cash contributed (₹)"] = summary_disp["Net cash contributed (₹)"].apply(_fmt_inr)
        st.dataframe(summary_disp, use_container_width=True)
        st.caption(
            "Net cash in/out per portfolio from the Cashflows sheet. TWR below is computed from the full "
            "daily ledger (equity + cash combined) — deposits/withdrawals are still excluded from the return "
            "itself, but idle cash and its timing are reflected in the portfolio's value."
        )

    st.subheader("Cash balance (latest)")
    if ledgers:
        cash_rows = {p: led["cash"].iloc[-1] for p, led in ledgers.items() if not led.empty}
        cash_disp = pd.Series(cash_rows, dtype="float64").sort_values(ascending=False).to_frame("Cash (₹)")
        cash_disp["Cash (₹)"] = cash_disp["Cash (₹)"].apply(_fmt_inr)
        st.dataframe(cash_disp, use_container_width=True)
    else:
        st.info("No ledger data available.")

    st.subheader("Data quality")
    if dl.failed_tickers:
        st.warning(f"Price download failed for {len(dl.failed_tickers)} tickers.")
        st.code("\n".join(dl.failed_tickers))
        if dl.failure_reasons:
            with st.expander("Why did these fail? (actual errors)"):
                reason_lines = [f"{t}: {dl.failure_reasons[t]}" for t in dl.failed_tickers if t in dl.failure_reasons]
                st.code("\n".join(reason_lines))
                st.caption(
                    "If the same error repeats across many tickers (especially large, liquid ones like "
                    "RELIANCE/HDFCBANK/SBIN), that's a systemic issue — a rate limit/IP block, an outdated "
                    "yfinance version, or a Yahoo Finance API change — not 40 individually bad tickers. "
                    "Worth checking the yfinance version pinned in requirements.txt and doing a full app "
                    "reboot (not just a rerun) on Streamlit Cloud so it reinstalls dependencies."
                )
    if missing_tickers:
        st.warning(f"Holdings contain {len(missing_tickers)} tickers missing from downloaded prices.")
        st.code("\n".join(sorted(missing_tickers)))
        with st.expander("Raw ticker text from your sheet (for these tickers)"):
            st.caption(
                "This is the literal text from the Ticker column before any processing — if a ticker "
                "shows something unexpected here (extra whitespace, an exchange prefix like 'NSE:', a "
                "different case, an already-present suffix in an odd form), that's the actual bug to fix, "
                "not the download logic."
            )
            raw_map = (
                df[df["ticker_yf"].isin(missing_tickers)][["ticker_yf", "ticker"]]
                .drop_duplicates()
                .rename(columns={"ticker_yf": "Resolved symbol", "ticker": "Raw text from sheet"})
                .sort_values("Resolved symbol")
            )
            st.dataframe(raw_map, use_container_width=True, hide_index=True)
        with st.expander("The actual download error for each missing ticker", expanded=True):
            st.caption(
                "The real exception (or 'no rows returned') from the last download attempt, per ticker. "
                "If the SAME error text repeats across most/all of these, that's a single systemic cause "
                "(e.g. a specific error class), not 40 independent bad tickers."
            )
            if dl.failure_reasons:
                reason_rows = [
                    {"Resolved symbol": t, "Last error": dl.failure_reasons[t]}
                    for t in sorted(missing_tickers) if t in dl.failure_reasons
                ]
                if reason_rows:
                    st.dataframe(pd.DataFrame(reason_rows), use_container_width=True, hide_index=True)
                else:
                    st.info("No failure reasons recorded for these tickers — see raw output below instead.")
            else:
                st.info("No failure reasons were captured at all this run.")
    if not backfill_summary.empty:
        st.info(
            f"{len(backfill_summary)} position(s) had a SELL with no matching BUY on record (owned before "
            "the tracked history began, but not entered in 'Initial Holdings'). An opening position was "
            "inferred automatically — sized to exactly cover the SELL, priced at the SELL's own price (so "
            "it assumes zero gain/loss on the untracked portion rather than guessing a cost basis) — and "
            "folded into every calculation below. For a precise cost basis, add these to 'Initial Holdings' "
            "with the real purchase price."
        )
        st.dataframe(backfill_summary, use_container_width=True, hide_index=True)
    if not oversold_warnings.empty:
        st.error(
            f"{len(oversold_warnings)} SELL transaction(s) still exceed recorded holdings even after "
            "auto-backfilling — this shouldn't happen and points to a bug in the backfill logic itself, "
            "not a data gap. Treat numbers for these tickers as unreliable."
        )
        warn_disp = oversold_warnings.copy()
        warn_disp["date"] = pd.to_datetime(warn_disp["date"]).dt.date.astype(str)
        st.dataframe(warn_disp, use_container_width=True)

    negative_cash_rows = []
    for _p, _led in ledgers.items():
        if _led.empty:
            continue
        _min_cash = float(_led["cash"].min())
        if _min_cash < -1.0:
            negative_cash_rows.append(
                {"portfolio": _p, "min_cash": round(_min_cash, 2), "date": _led["cash"].idxmin().date().isoformat()}
            )
    if negative_cash_rows:
        st.error(
            f"{len(negative_cash_rows)} portfolio(s) show negative cash at some point in the daily ledger. "
            "This happens when a BUY isn't matched by a recorded Cashflows deposit — for a portfolio with "
            "'Initial Holdings', that position is treated as already-funded (no cash debit at StartDate), so "
            "any *subsequent* real purchase needs its own deposit on record. A negative cash balance "
            "understates Portfolio Value and will distort TWR, Sharpe, Sortino, and every other ratio "
            "downstream for that portfolio — this is very likely why the ratios look off."
        )
        st.dataframe(pd.DataFrame(negative_cash_rows), use_container_width=True)


with client_tab:
    portfolios = portfolios_all
    if not portfolios:
        st.error("No portfolios found in the uploaded data.")
        st.stop()

    portfolio = st.selectbox("Client / portfolio", portfolios)

    # Transactions filtered to selected client
    st.subheader("Transactions")
    client_txns = df[df["portfolio"].astype(str).str.strip() == str(portfolio).strip()]
    st.dataframe(client_txns.sort_values("date", ascending=False).reset_index(drop=True), use_container_width=True)

    st.subheader("Cash flows")
    client_cash = cashflows_df[cashflows_df["portfolio"].astype(str).str.strip() == str(portfolio).strip()]
    if client_cash.empty:
        st.info("No cashflow records for this client.")
    else:
        cash_disp = client_cash.sort_values("date", ascending=False).reset_index(drop=True).copy()
        cash_disp["amount"] = cash_disp["amount"].apply(_fmt_inr)
        st.dataframe(cash_disp, use_container_width=True)

    pr = twr_df[portfolio].dropna() if portfolio in twr_df.columns else pd.Series(dtype=float)
    corr = None
    client_ledger = ledgers.get(str(portfolio), pd.DataFrame())

    # ---- Performance Summary / Performance Table ----------------------------
    # Matches the reference PortfolioTWR.py's own dashboard layout exactly,
    # computed via the same functions (build_daily_ledger, compute_ledger_twr,
    # annualise_twr_simple, benchmark_existing/benchmark_replication) so the
    # numbers here are the ones to trust and compare against.
    st.subheader("Performance Summary")
    if not client_ledger.empty:
        cum_twr = float((1 + pr).prod() - 1) if not pr.empty else np.nan
        days_elapsed = (client_ledger.index.max() - client_ledger.index.min()).days + 1
        ann_twr = annualise_twr_simple(cum_twr, days_elapsed)

        s1, s2, s3, s4 = st.columns(4)
        s1.metric("Portfolio TWR", _fmt_pct(cum_twr))
        s2.metric("Annualised TWR", _fmt_pct(ann_twr))
        s3.metric("Current Portfolio Value", _fmt_inr(client_ledger["portfolio_value"].iloc[-1]))
        s4.metric("Cash Balance", _fmt_inr(client_ledger["cash"].iloc[-1]))

        # Determine whether this portfolio has existing (Initial Holdings)
        # opening, to pick the matching benchmark simulation method.
        client_key = str(portfolio).strip()
        client_txns_all = df[df["portfolio"].astype(str).str.strip() == client_key]
        client_cf_all = cashflows_df[cashflows_df["portfolio"].astype(str).str.strip() == client_key]
        client_opening_rows = (
            client_txns_all[client_txns_all.get("is_opening", False) == True]  # noqa: E712
            if "is_opening" in client_txns_all.columns
            else pd.DataFrame()
        )
        client_has_existing = False
        opening_total_val = 0.0
        if not client_opening_rows.empty and "opening_market_value" in client_opening_rows.columns:
            mv_sum = pd.to_numeric(client_opening_rows["opening_market_value"], errors="coerce").sum()
            if pd.notna(mv_sum) and mv_sum > 0:
                client_has_existing = True
                opening_total_val = float(mv_sum)

        perf_table = pd.DataFrame(
            {"Portfolio": [cum_twr, ann_twr, np.nan, np.nan]},
            index=["Actual TWR", "Annualised TWR", "Alpha vs Nifty 50", "Alpha vs Nifty 500"],
        )
        for name in ["Nifty 50", "Nifty 500"]:
            if name not in benchmarks:
                continue
            bench_series = benchmarks[name].reindex(client_ledger.index).ffill()
            if client_has_existing:
                bench_ledger = benchmark_existing(
                    client_cf_all, bench_series, opening_total_val, client_ledger.index.min()
                )
            else:
                bench_ledger = benchmark_replication(client_cf_all, bench_series)
            bench_twr_series = compute_ledger_twr(bench_ledger)
            bench_cum = float((1 + bench_twr_series).prod() - 1) if not bench_twr_series.empty else np.nan
            bench_ann = annualise_twr_simple(bench_cum, days_elapsed)
            col = f"{name} TRI"
            perf_table[col] = [bench_cum, bench_ann, np.nan, np.nan]
            alpha = cum_twr - bench_cum if pd.notna(cum_twr) and pd.notna(bench_cum) else np.nan
            perf_table.loc[f"Alpha vs {name}", "Portfolio"] = alpha

        perf_table_fmt = perf_table.apply(lambda col: col.apply(lambda v: "N/A" if pd.isna(v) else f"{v*100:.2f}%"))
        st.subheader("Performance Table")
        st.dataframe(perf_table_fmt, use_container_width=True)

        st.subheader("Trailing Period Returns")
        st.caption(
            "Portfolio TWR vs. Nifty 500's own price return, over standard trailing windows, with Alpha "
            "(Portfolio minus Nifty 500) for each period — a blank cell means there isn't enough history "
            "yet to cover that period."
        )
        trailing_bench_series = {
            "Nifty 500 TRI": benchmarks["Nifty 500"].reindex(client_ledger.index).ffill()
        } if "Nifty 500" in benchmarks else {}
        trailing_table = trailing_returns_table(pr, trailing_bench_series)
        if trailing_table.empty:
            st.info("Not enough data to compute trailing period returns.")
        else:
            trailing_fmt = trailing_table.apply(lambda col: col.apply(lambda v: "" if pd.isna(v) else f"{v*100:.2f}%"))
            st.dataframe(trailing_fmt, use_container_width=True)
    else:
        st.info("Not enough data to compute a performance summary for this client.")

    st.subheader("Extended statistics")
    st.caption("Supplementary risk-adjusted metrics not present in the reference app, computed from the same TWR series above.")
    client_metrics = portfolio_metrics(pr, benchmarks=benchmarks, rf_annual=rf_annual)

    v1, v2 = st.columns(2)
    if not client_ledger.empty:
        v1.metric("Current portfolio value", _fmt_inr(client_ledger["portfolio_value"].iloc[-1]))
        v2.metric("Cash balance", _fmt_inr(client_ledger["cash"].iloc[-1]))
    else:
        v1.metric("Current portfolio value", "N/A")
        v2.metric("Cash balance", "N/A")

    k1, k2, k3, k4 = st.columns(4)
    k1.metric("Period return (TWR)", _fmt_pct(client_metrics.get("Period Return (TWR)")))
    k2.metric("Period volatility (TWR)", _fmt_pct(client_metrics.get("Period Volatility (TWR)")))
    k3.metric("Sharpe", _fmt_float(client_metrics.get("Sharpe")))
    k4.metric("Max drawdown", _fmt_pct(client_metrics.get("Max Drawdown")))
    m1, m2 = st.columns(2)
    m1.metric("Sortino", _fmt_float(client_metrics.get("Sortino")))
    m2.metric("99% CVaR (Daily)", _fmt_pct(client_metrics.get("CVaR 99% (Daily)")))

    if benchmarks:
        st.caption("vs. benchmark")
        bench_rows = []
        for name in benchmarks.keys():
            bench_rows.append(
                {
                    "Benchmark": name,
                    "Beta": _fmt_float(client_metrics.get(f"Beta ({name})")),
                    "Jensen Alpha (TWR trailing period)": _fmt_pct(client_metrics.get(f"Jensen Alpha ({name})")),
                    "Treynor": _fmt_float(client_metrics.get(f"Treynor ({name})")),
                }
            )
        st.dataframe(pd.DataFrame(bench_rows).set_index("Benchmark"), use_container_width=True)

    growth = cumulative_growth_from_twr(pr)

    c1, c2 = st.columns([2, 1])
    with c1:
        st.subheader("Performance vs. Nifty 50 & Nifty 500")
        if growth.empty:
            st.info("Not enough return history for a performance chart.")
        else:
            # Rebase each benchmark to 1.0 on the portfolio's own start date so
            # all lines are directly comparable growth-of-₹1 curves over the
            # same window, regardless of the benchmark's own longer history.
            comparison = pd.DataFrame({str(portfolio): growth})
            for name, bench_prices in benchmarks.items():
                bench_aligned = bench_prices.reindex(growth.index).ffill()
                if bench_aligned.dropna().empty:
                    continue
                first_valid = bench_aligned.dropna().iloc[0]
                if first_valid and first_valid > 0:
                    comparison[name] = bench_aligned / first_valid
            st.line_chart(comparison, height=260)
    with c2:
        st.subheader("Drawdown")
        dd = drawdown_series(growth)
        if dd.empty:
            st.info("Not enough return history for drawdown.")
        else:
            import matplotlib.pyplot as plt
            import matplotlib.ticker as mticker
            fig_dd, ax_dd = plt.subplots(figsize=(6, 2.8))
            ax_dd.fill_between(dd.index, dd.values * 100, 0, color="#AA1437", alpha=0.6)
            ax_dd.plot(dd.index, dd.values * 100, color="#922B21", linewidth=0.8)
            ax_dd.yaxis.set_major_formatter(mticker.FormatStrFormatter("%.1f%%"))
            ax_dd.set_ylabel("Drawdown (%)")
            ax_dd.set_title(f"Drawdown — {portfolio}", fontsize=9)
            ax_dd.grid(False)
            # Ensure y-axis shows the actual range (negative values)
            dd_min = float(dd.min()) * 100
            ax_dd.set_ylim(min(dd_min * 1.15, dd_min - 1), 1)
            fig_dd.tight_layout()
            st.pyplot(fig_dd, use_container_width=True)
            plt.close(fig_dd)

    st.subheader("Returns (daily)")
    st.line_chart(pr, height=220)

    # Normalised positions for this portfolio — used in both Holdings and Top Movers
    pos_copy = positions.copy()
    pos_copy["portfolio"] = pos_copy["portfolio"].astype(str).str.strip()
    port_pos = pos_copy[pos_copy["portfolio"] == str(portfolio).strip()]

    st.subheader("Holdings (latest)")
    holdings = daily_positions.get(str(portfolio), pd.DataFrame())
    if holdings.empty:
        st.info("No holdings available for this portfolio.")
    else:
        latest = holdings.iloc[-1]
        latest = latest[latest > 0].sort_values(ascending=False).to_frame(name="Qty")
        if "avg_buy_price" in port_pos.columns and not port_pos.empty:
            cur_pos = port_pos[["ticker_yf", "avg_buy_price"]].set_index("ticker_yf")
            cur_pos = cur_pos.rename(columns={"avg_buy_price": "Avg Buy Price (₹)"})
            latest = latest.join(cur_pos, how="left")
        st.dataframe(latest, use_container_width=True)

    st.subheader("Market cap allocation")
    cap_split = pd.Series(dtype=float)
    if not holdings.empty:
        held_tickers = tuple(sorted(latest.index.tolist())) if not latest.empty else ()
        if held_tickers:
            with st.spinner("Fetching market caps…"):
                mkt_caps = _fetch_market_caps_cached(held_tickers)
            cap_split = cap_split_weights(holdings, prices, mkt_caps)
        if cap_split.empty:
            st.info("Market cap data unavailable for this portfolio's holdings.")
        else:
            # matplotlib, not plotly — plotly isn't a guaranteed dependency on
            # every deployment and every other chart in this app already uses
            # matplotlib, so this stays consistent and doesn't add a new
            # import that can break the app if it's missing at deploy time.
            import matplotlib.pyplot as plt

            cap_color_map = {"Large Cap": "#290084", "Mid Cap": "#F6BF02", "Small Cap": "#1D9951", "ETF": "#717D7E", "REIT": "#AA1437", "Unknown": "#B0B0B0"}
            colors_ordered = [cap_color_map.get(k, "#B0B0B0") for k in cap_split.index]
            fig_cap, ax_cap = plt.subplots(figsize=(5, 5))
            _, _, autotexts_cap = ax_cap.pie(
                cap_split.values, labels=cap_split.index, autopct="%1.1f%%", startangle=90,
                colors=colors_ordered, textprops={"fontsize": 9},
            )
            for txt, color in zip(autotexts_cap, colors_ordered):
                txt.set_color(_contrast_text_color(color))
            fig_cap.tight_layout()
            left_col, center_col, right_col = st.columns([1, 2, 1])
            with center_col:
                st.pyplot(fig_cap, use_container_width=True)
            plt.close(fig_cap)
    else:
        st.info("No holdings available for this portfolio.")

    st.subheader("Sector exposure")
    # Sector classification is purely from the SECTOR_MAP dict in code —
    # no yfinance lookup. An unlisted ticker shows up as "Unknown"; add it
    # to SECTOR_MAP directly in the code to fix that, the same way the
    # existing entries were curated by hand.
    effective_sector_map = SECTOR_MAP
    if not holdings.empty and held_tickers:
        sector_split, sector_counts = sector_exposure_weights_and_counts(holdings, prices, effective_sector_map)
        if sector_split.empty:
            st.info("Sector data unavailable for this portfolio's holdings.")
        else:
            fig_sec, colors_hex_sec = sector_pie_figure_clean(sector_split)
            pie_col, legend_col = st.columns([1, 1])
            with pie_col:
                st.pyplot(fig_sec, use_container_width=True)
            plt.close(fig_sec)
            with legend_col:
                legend_rows = "".join(
                    f'<tr><td style="padding:3px 8px;">'
                    f'<span style="display:inline-block;width:12px;height:12px;background:{c};'
                    f'border-radius:2px;"></span></td>'
                    f'<td style="padding:3px 8px;">{sector}</td>'
                    f'<td style="padding:3px 8px;text-align:right;">{weight*100:.1f}%</td>'
                    f'<td style="padding:3px 8px;text-align:right;">{int(sector_counts.get(sector, 0))}</td></tr>'
                    for c, (sector, weight) in zip(colors_hex_sec, sector_split.items())
                )
                st.markdown(
                    f"""
                    <table style="width:100%; border-collapse:collapse; font-size:14px;">
                        <thead>
                            <tr style="border-bottom:2px solid #290084;">
                                <th></th>
                                <th style="text-align:left; padding:3px 8px;">Sector</th>
                                <th style="text-align:right; padding:3px 8px;">%</th>
                                <th style="text-align:right; padding:3px 8px;">Stocks</th>
                            </tr>
                        </thead>
                        <tbody>{legend_rows}</tbody>
                    </table>
                    """,
                    unsafe_allow_html=True,
                )
            if (sector_split.index == "Unknown").any():
                unclassified_tickers = sorted(
                    t for t in held_tickers if t not in effective_sector_map
                )
                st.caption(
                    f"{sector_split.get('Unknown', 0)*100:.1f}% not in SECTOR_MAP yet: "
                    + ", ".join(t.replace(".NS", "").replace(".BO", "") for t in unclassified_tickers)
                    + ". Add these to SECTOR_MAP in the code to classify them."
                )
    else:
        st.info("No holdings available for this portfolio.")

    st.subheader("Download")
    report_choice = st.selectbox(
        "Choose report",
        [
            "PDF report",
            "PowerPoint report (editable)",
            "Cumulative growth (CSV)",
            "Daily returns (CSV)",
            "Holdings (latest, CSV)",
            "Transactions (this client, CSV)",
        ],
    )

    report_date = datetime.now().date()
    stats_row = client_metrics
    benchmark_returns = {name: s.pct_change() for name, s in benchmarks.items()} if benchmarks else None

    st.subheader("Attribution & risk")
    attrib = None
    rc = None
    if not holdings.empty:
        attrib = performance_attribution(holdings, prices)
        rc = risk_contribution(holdings, prices)

        st.caption("Contribution to return, since inception — biggest positive to biggest negative contributor")
        if attrib is not None and not attrib.empty:
            contrib_df = attrib.sort_values(ascending=False).to_frame("contribution")
            total_return_sum = contrib_df["contribution"].sum()
            contrib_df.index = contrib_df.index.str.replace(".NS", "", regex=False).str.replace(".BO", "", regex=False)
            contrib_df.index.name = "Ticker"
            contrib_df["Contribution to Return"] = contrib_df["contribution"].apply(_fmt_pct_signed)
            # Normalized to sum to 100% across all tickers — "this stock was
            # responsible for X% of the whole return", alongside the existing
            # percentage-point figure ("this stock added X points").
            if total_return_sum and abs(total_return_sum) > 1e-9:
                contrib_df["% of Total Return"] = (contrib_df["contribution"] / total_return_sum).apply(_fmt_pct_signed)
            else:
                contrib_df["% of Total Return"] = "N/A"
            contrib_df = contrib_df[["Contribution to Return", "% of Total Return"]]
            st.dataframe(
                contrib_df,
                use_container_width=True,
                height=min(35 * (len(contrib_df) + 1), 480),
            )
        else:
            st.info("Not enough data for attribution.")

        st.caption("Risk contribution — contribution to total portfolio volatility, and % of total portfolio risk")
        if rc is not None and not rc.empty:
            port_vol = annualized_volatility(pr) if not pr.empty else np.nan
            rc_df = rc.sort_values(ascending=False).to_frame("share")
            rc_df.index = rc_df.index.str.replace(".NS", "", regex=False).str.replace(".BO", "", regex=False)
            rc_df.index.name = "Ticker"
            # Euler decomposition: share_i * portfolio_vol sums exactly back
            # to portfolio_vol across all tickers — the "current numbers"
            # (contribution in actual volatility percentage-points) alongside
            # the normalized share (already summed to 100% by risk_contribution).
            if pd.notna(port_vol):
                rc_df["Contribution to Volatility"] = (rc_df["share"] * port_vol).apply(_fmt_pct_signed)
            else:
                rc_df["Contribution to Volatility"] = "N/A"
            rc_df["% of Total Risk"] = rc_df["share"].apply(_fmt_pct_signed)
            rc_df = rc_df[["Contribution to Volatility", "% of Total Risk"]]
            st.dataframe(rc_df, use_container_width=True, height=min(35 * (len(rc_df) + 1), 480))
        else:
            currently_held_ct = int((holdings.iloc[-1] > 1e-9).sum()) if not holdings.empty else 0
            if currently_held_ct < 2:
                st.info(
                    "Only one current holding in this portfolio — risk contribution needs at least two "
                    "to show how risk is split between them."
                )
            else:
                st.info(
                    "Not enough overlapping price history between the current holdings to compute a "
                    "reliable covariance (e.g. a recently-listed stock with a short price history)."
                )

        if not pr.empty:
            recon = attribution_reconciliation(holdings, prices, pr)
            gap = recon.get("gap")
            if pd.notna(gap):
                st.caption(
                    f"Sum of ticker contributions: {_fmt_pct(recon['contribution_total'])} vs. actual cumulative "
                    f"return: {_fmt_pct(recon['actual_cumulative_return'])} (gap: {_fmt_pct(gap)}, "
                    "expected from daily-rebalanced attribution vs. compounded return)."
                )

        st.subheader("Risk vs. return contribution by holding")
        st.caption(
            "Not each stock's own standalone return/volatility — this plots what % of the portfolio's TOTAL "
            "return and TOTAL risk (both above, both summing to 100%) each stock is actually responsible for. "
            "Dashed lines mark the median across this portfolio's own holdings, splitting them into quadrants "
            "relative to each other; the solid lines at zero separate positive from negative contributors."
        )
        rr_df = stock_risk_return(
            attrib if attrib is not None else pd.Series(dtype=float),
            rc if rc is not None else pd.Series(dtype=float),
        )
        if rr_df.empty:
            st.info("Not enough attribution/risk data to plot risk vs. return contribution for this portfolio.")
        else:
            fig_rr = risk_return_scatter_figure(rr_df)
            if fig_rr is not None:
                st.pyplot(fig_rr, use_container_width=True)
                plt.close(fig_rr)
            rr_labeled = risk_return_quadrants(rr_df).sort_values("return_contribution", ascending=False)
            rr_disp = rr_labeled.set_index("ticker")[["return_contribution", "risk_contribution", "quadrant"]].rename(
                columns={
                    "return_contribution": "% of Total Return",
                    "risk_contribution": "% of Total Risk",
                    "quadrant": "Quadrant",
                }
            )
            rr_disp["% of Total Return"] = rr_disp["% of Total Return"].apply(_fmt_pct_signed)
            rr_disp["% of Total Risk"] = rr_disp["% of Total Risk"].apply(_fmt_pct_signed)
            st.dataframe(rr_disp, use_container_width=True)

    st.subheader("Return distribution (daily)")
    if not pr.empty and len(pr) >= 10:
        import matplotlib.pyplot as plt
        import matplotlib.ticker as mticker
        cutoff_99 = float(pr.quantile(0.01))
        tail_99   = pr[pr <= cutoff_99]
        cvar99_daily = float(-tail_99.mean()) if not tail_99.empty else float(-cutoff_99)
        fig_hist, ax_hist = plt.subplots(figsize=(9, 5.2))
        counts, bin_edges, _ = ax_hist.hist(
            pr.dropna().values, bins=40, color="#F6BF02", alpha=0.85, edgecolor="white", linewidth=0.3,
            label="Daily returns",
        )
        mu_h, sigma_h = float(pr.mean()), float(pr.std(ddof=1))
        if sigma_h > 0 and len(bin_edges) > 1:
            bin_width_h = bin_edges[1] - bin_edges[0]
            x_curve = np.linspace(pr.min(), pr.max(), 200)
            pdf = (1.0 / (sigma_h * np.sqrt(2 * np.pi))) * np.exp(-0.5 * ((x_curve - mu_h) / sigma_h) ** 2)
            y_curve = pdf * len(pr.dropna()) * bin_width_h
            ax_hist.plot(x_curve, y_curve, color="#290084", linewidth=2, label="Fitted normal distribution")
        ax_hist.axvline(cutoff_99, color="#AA1437", linewidth=1.5, linestyle="--",
                        label=f"99% CVaR: {cvar99_daily*100:.2f}%  (tail avg beyond {cutoff_99*100:.2f}%)")
        ax_hist.fill_betweenx([0, ax_hist.get_ylim()[1] or 1],
                              pr.min(), cutoff_99,
                              alpha=0.12, color="#AA1437", label="_nolegend_")
        ax_hist.set_xlabel("Daily return")
        ax_hist.set_ylabel("Frequency")
        ax_hist.set_title(f"Distribution of daily returns — {portfolio}", fontsize=10)
        ax_hist.xaxis.set_major_formatter(mticker.PercentFormatter(xmax=1, decimals=1))
        ax_hist.legend(fontsize=8)
        ax_hist.grid(False)
        fig_hist.tight_layout()
        st.pyplot(fig_hist, use_container_width=True)
        plt.close(fig_hist)

        st.caption("Distribution statistics")
        stats_tbl = pd.DataFrame(
            {
                "Value": {
                    "Mean (daily)": f"{float(pr.mean())*100:+.3f}%",
                    "Std. Dev. (daily)": f"{float(pr.std(ddof=1))*100:.3f}%",
                    "Skewness": f"{float(pr.skew()):+.2f}",
                    "Kurtosis (excess)": f"{float(pr.kurt()):+.2f}",
                    "Min (daily)": f"{float(pr.min())*100:+.2f}%",
                    "Max (daily)": f"{float(pr.max())*100:+.2f}%",
                }
            }
        )
        st.dataframe(stats_tbl, use_container_width=True)
    else:
        st.info("Not enough return history to display distribution.")

    st.subheader("Stock correlation (this client)")
    currently_held = set(
        pos_copy.loc[pos_copy["portfolio"] == str(portfolio).strip(), "ticker_yf"].tolist()
    )
    tickers_in_portfolio = [t for t in currently_held if t in prices.columns]
    if len(tickers_in_portfolio) >= 2:
        stock_rets = prices[tickers_in_portfolio].pct_change().dropna(how="all")
        stock_rets = stock_rets.dropna(axis=1, how="all")
        if stock_rets.shape[1] >= 2:
            corr = stock_rets.corr().round(2)
            st.dataframe(corr, use_container_width=True)
            _heatmap(corr, title=f"Stock correlation heatmap — {portfolio}")
            corr_summary = correlation_summary(
                corr, pr, benchmarks["Nifty 500"].pct_change() if "Nifty 500" in benchmarks else None
            )
            cs1, cs2 = st.columns(2)
            cs1.metric("Avg. pairwise correlation (holdings)", _fmt_float(corr_summary["avg_pairwise"]))
            cs2.metric("Portfolio correlation vs. Nifty 500", _fmt_float(corr_summary["vs_benchmark"]))
        else:
            st.info("Not enough price history for a stock correlation matrix.")
    else:
        st.info("Need at least 2 stocks in this portfolio to show correlation.")

    if report_choice == "PDF report":
        if holdings.empty or pr.empty:
            st.warning("Not enough data to generate a PDF report for this client.")
        else:
            try:
                pdf_bytes = generate_portfolio_pdf_bytes(
                    portfolio_name=str(portfolio),
                    report_date=report_date,
                    returns=pr,
                    holdings=holdings,
                    prices=prices,
                    stats_row=stats_row,
                    benchmark_returns=benchmark_returns,
                    sector_map=effective_sector_map,
                    risk_contrib=rc if isinstance(rc, pd.Series) and not rc.empty else None,
                    attribution=attrib if isinstance(attrib, pd.Series) and not attrib.empty else None,
                    stock_corr=corr if "corr" in locals() else None,
                    returns_hist=pr if not pr.empty else None,
                    cash_balance=float(client_ledger["cash"].iloc[-1]) if not client_ledger.empty else None,
                    portfolio_value=float(client_ledger["portfolio_value"].iloc[-1]) if not client_ledger.empty else None,
                    cap_split=cap_split if isinstance(cap_split, pd.Series) and not cap_split.empty else None,
                    trailing_returns=trailing_table if "trailing_table" in locals() and isinstance(trailing_table, pd.DataFrame) and not trailing_table.empty else None,
                )
                st.download_button(
                    "Download PDF",
                    data=pdf_bytes,
                    file_name=f"Portfolio_{portfolio}.pdf",
                    mime="application/pdf",
                    key=f"dl_pdf_{portfolio}",
                )
            except Exception as e:
                import traceback
                st.error(f"Failed to build PDF: {e}")
                st.code(traceback.format_exc())

    elif report_choice == "PowerPoint report (editable)":
        if holdings.empty or pr.empty:
            st.warning("Not enough data to generate a PowerPoint report for this client.")
        else:
            try:
                pptx_bytes = generate_portfolio_pptx_bytes(
                    portfolio_name=str(portfolio),
                    report_date=report_date,
                    returns=pr,
                    holdings=holdings,
                    prices=prices,
                    stats_row=stats_row,
                    benchmark_returns=benchmark_returns,
                    sector_map=effective_sector_map,
                    risk_contrib=rc if isinstance(rc, pd.Series) and not rc.empty else None,
                    attribution=attrib if isinstance(attrib, pd.Series) and not attrib.empty else None,
                    stock_corr=corr if "corr" in locals() else None,
                    returns_hist=pr if not pr.empty else None,
                    cash_balance=float(client_ledger["cash"].iloc[-1]) if not client_ledger.empty else None,
                    portfolio_value=float(client_ledger["portfolio_value"].iloc[-1]) if not client_ledger.empty else None,
                    cap_split=cap_split if isinstance(cap_split, pd.Series) and not cap_split.empty else None,
                    trailing_returns=trailing_table if "trailing_table" in locals() and isinstance(trailing_table, pd.DataFrame) and not trailing_table.empty else None,
                )
                st.download_button(
                    "Download PowerPoint",
                    data=pptx_bytes,
                    file_name=f"Portfolio_{portfolio}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"dl_pptx_{portfolio}",
                )
                st.caption(
                    "Text and tables are fully editable in PowerPoint. Charts are embedded as images "
                    "(same as every BI tool's PPT export) — restyle the tables or add your own slides "
                    "around them freely."
                )
            except Exception as e:
                import traceback
                st.error(f"Failed to build PowerPoint: {e}")
                st.code(traceback.format_exc())

    elif report_choice == "Cumulative growth (CSV)":
        out = growth.to_frame(name="growth_of_1_rupee")
        st.download_button(
            "Download growth CSV",
            data=_df_to_csv_bytes(out),
            file_name=f"Growth_{portfolio}.csv",
            mime="text/csv",
            key=f"dl_growth_{portfolio}",
        )

    elif report_choice == "Daily returns (CSV)":
        out = pr.to_frame(name="returns")
        st.download_button(
            "Download returns CSV",
            data=_df_to_csv_bytes(out),
            file_name=f"Returns_{portfolio}.csv",
            mime="text/csv",
            key=f"dl_rets_{portfolio}",
        )

    elif report_choice == "Holdings (latest, CSV)":
        if holdings.empty:
            st.warning("No holdings for this portfolio.")
        else:
            out = holdings.iloc[-1].to_frame(name="qty")
            out = out[out["qty"] > 0].sort_values("qty", ascending=False)
            st.download_button(
                "Download holdings CSV",
                data=_df_to_csv_bytes(out),
                file_name=f"Holdings_{portfolio}.csv",
                mime="text/csv",
                key=f"dl_hold_{portfolio}",
            )

    else:  # Transactions
        out = df[df["portfolio"] == portfolio].copy()
        st.download_button(
            "Download transactions CSV",
            data=_df_to_csv_bytes(out),
            file_name=f"Transactions_{portfolio}.csv",
            mime="text/csv",
            key=f"dl_txn_{portfolio}",
        )
